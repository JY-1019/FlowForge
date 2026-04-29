"""Built-in tools for dynamic FlowForge agents.

Two categories of tools are provided:

1. **Shell / runtime tools** — project inspection, test running, dependency
   checking.  Gated by ``DynamicRunOptions.allowed_shell_modes``.
2. **General-purpose utility tools** — ``web_fetch_url``, ``json_select_fields``,
   ``files_read_text``, ``files_write_text``, ``files_list_dir``.  Always
   included when ``include_builtin_tools=True`` (the default).  These are
   pure-Python tools that never shell out.
"""
from __future__ import annotations

import shlex
import socket
import subprocess
import time
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

from flowforge.types import DynamicRunOptions, FunctionTool

_CONTROL_TOKENS = (";", "&&", "||", "|", ">", "<", "$(", "`")
_DEFAULT_MAX_OUTPUT_CHARS = 4000


def create_builtin_tool_pack(options: DynamicRunOptions) -> list[FunctionTool]:
    """Return builtin tools enabled by *options*.

    Shell tools are gated by ``options.allowed_shell_modes``.  General-purpose
    utility tools (web, json, files) are always included — they are safe,
    sandboxed, and high-reuse.
    """
    tools: list[FunctionTool] = []
    modes = set(options.allowed_shell_modes)

    # ── Always-on utility tools ───────────────────────────────────────
    tools.append(FunctionTool(
        func=_make_import_check_tool(),
        name="python_import_check",
        description=(
            "Check whether Python modules can be imported and report their "
            "installed package versions when available."
        ),
    ))
    tools.append(FunctionTool(
        func=_make_mcp_start_tool(options),
        name="mcp_start_server",
        description=(
            "Start an MCP server command declared in DynamicRunOptions."
        ),
    ))
    tools.append(FunctionTool(
        func=_make_mcp_register_tool(options),
        name="mcp_register_server",
        description=(
            "Register one or more MCP tool names from a declared MCP server "
            "URL into the current FlowForge run so later steps can use them "
            "with <tool_name> in ctx.call_llm(). If a server command is "
            "declared and the endpoint is not reachable, this tool starts it "
            "automatically by default."
        ),
    ))

    tools.append(FunctionTool(
        func=_make_pip_install_tool(options),
        name="pip_install",
        description=(
            "Install one or more Python packages using pip. Accepts a "
            "comma-separated or space-separated list of package names "
            "(e.g. 'python-pptx, httpx'). Installation is gated by the "
            "agent's DependencyPolicy."
        ),
    ))

    # General-purpose utility tools.
    tools.extend(_create_utility_tools(options))

    # Document processing tools (PDF, PPTX, CSV, DOCX, Markdown, Chart, Image).
    tools.extend(_create_document_tools(options))

    # Claude Code skill invocation tool.
    tools.append(FunctionTool(
        func=_make_claude_skill_tool(options),
        name="claude_skill",
        description=(
            "Invoke a Claude Code skill (slash command) and return the "
            "result. Pass the skill name (e.g. 'commit', 'review-pr') "
            "and a prompt string. Requires the 'claude' CLI in PATH."
        ),
    ))

    # ── Shell tools (mode-gated) ──────────────────────────────────────
    if "readonly" in modes:
        tools.append(FunctionTool(
            func=_make_shell_tool(options, "readonly"),
            name="shell_readonly",
            description=(
                "Run a read-only shell command inside the project. Useful for "
                "pwd, ls, find, rg, cat, sed, head, tail, wc, and git inspect."
            ),
        ))
    if "project_exec" in modes:
        tools.append(FunctionTool(
            func=_make_shell_tool(options, "project_exec"),
            name="shell_project_exec",
            description=(
                "Run a project command such as python -m pytest, pytest, uv, "
                "npm, pnpm, or yarn inside the project."
            ),
        ))
    if "workspace_write" in modes:
        tools.append(FunctionTool(
            func=_make_shell_tool(options, "workspace_write"),
            name="shell_workspace_write",
            description=(
                "Run a limited workspace-writing shell command inside the "
                "project, such as mkdir, touch, cp, or mv."
            ),
        ))
    if "install_dependency" in modes:
        tools.append(FunctionTool(
            func=_make_shell_tool(options, "install_dependency"),
            name="shell_install_dependency",
            description=(
                "Install a dependency with pip, uv, npm, pnpm, or yarn when "
                "the dynamic dependency policy allows installation."
            ),
        ))

    return tools


# ---------------------------------------------------------------------------
# General-purpose utility tools (always included)
# ---------------------------------------------------------------------------

def _create_utility_tools(options: DynamicRunOptions) -> list[FunctionTool]:
    """Create pure-Python utility tools for web, JSON, and file operations."""
    project_root = Path(options.project_root or Path.cwd()).expanduser().resolve()
    max_output = max(500, options.shell_output_max_chars or _DEFAULT_MAX_OUTPUT_CHARS)

    return [
        FunctionTool(
            func=_make_web_fetch_url_tool(max_output),
            name="web_fetch_url",
            description=(
                "Fetch a URL via HTTP GET. Returns dict {ok, status, "
                "content_type, body, truncated, url}; use result['body'] "
                "for the response text. Pass max_chars to request a larger "
                "or smaller body slice; max_chars<=0 returns the full body."
            ),
        ),
        FunctionTool(
            func=_make_json_select_fields_tool(),
            name="json_select_fields",
            description=(
                "Select specific top-level fields from a JSON string or dict. "
                "Returns a new dict containing only the requested keys."
            ),
        ),
        FunctionTool(
            func=_make_files_read_text_tool(project_root, max_output),
            name="files_read_text",
            description=(
                "Read a text file in the project (path relative to project "
                "root). Returns dict {ok, content, truncated, path, size} "
                "or {ok: False, error}; use result['content'] for the text."
            ),
        ),
        FunctionTool(
            func=_make_files_write_text_tool(project_root),
            name="files_write_text",
            description=(
                "Write text to a file in the project (path relative to "
                "project root or absolute under it). Creates parent dirs. "
                "Returns dict {ok, path, size} or {ok: False, error}."
            ),
        ),
        FunctionTool(
            func=_make_files_list_dir_tool(project_root),
            name="files_list_dir",
            description=(
                "List files and directories at a path inside the project. "
                "Path must be relative to the project root."
            ),
        ),
    ]


def _make_web_fetch_url_tool(default_max_chars: int):
    def _tool(
        url: str,
        timeout_seconds: int = 30,
        max_chars: int | None = None,
    ) -> dict[str, Any]:
        # Use httpx instead of urllib so we can:
        #   * disable TLS verification behind corporate proxies that
        #     re-sign certs (matches the LLMConfig.verify_ssl=False
        #     pattern used elsewhere);
        #   * follow redirects automatically;
        #   * honour HTTP_PROXY / HTTPS_PROXY environment variables.
        # urllib silently returns an empty body in some proxy setups,
        # which is the failure mode we hit in production.
        try:
            import httpx
        except ImportError as e:
            return {
                "ok": False,
                "status": 0,
                "error": f"httpx not available: {e}",
                "url": url,
            }

        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/120.0 Safari/537.36 FlowForge/1.0"
            ),
            "Accept": (
                "text/html,application/xhtml+xml,application/xml;q=0.9,"
                "application/json;q=0.8,*/*;q=0.7"
            ),
            "Accept-Language": "en-US,en;q=0.9,ko;q=0.8",
        }

        try:
            with httpx.Client(
                verify=False,
                follow_redirects=True,
                timeout=timeout_seconds,
                trust_env=True,
            ) as client:
                resp = client.get(url, headers=headers)
            body = resp.text or ""
            content_type = resp.headers.get("Content-Type", "")
            limit = default_max_chars if max_chars is None else int(max_chars)
            if limit <= 0:
                returned_body = body
                truncated = False
            else:
                returned_body = body[:limit]
                truncated = len(body) > limit
            return {
                "ok": True,
                "status": resp.status_code,
                "content_type": content_type,
                "body": returned_body,
                "truncated": truncated,
                "url": str(resp.url),
            }
        except httpx.HTTPStatusError as e:
            return {
                "ok": False,
                "status": e.response.status_code,
                "error": str(e),
                "url": url,
            }
        except Exception as e:
            return {
                "ok": False,
                "status": 0,
                "error": str(e),
                "url": url,
            }

    _tool.__name__ = "builtin_web_fetch_url"
    return _tool


def _make_json_select_fields_tool():
    import json as _json

    def _tool(data: str, fields: str) -> dict[str, Any]:
        try:
            if isinstance(data, str):
                parsed = _json.loads(data)
            else:
                parsed = data
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid JSON: {e}"}

        if not isinstance(parsed, dict):
            return {"ok": False, "error": "Input is not a JSON object"}

        field_list = [f.strip() for f in fields.split(",") if f.strip()]
        selected = {k: parsed[k] for k in field_list if k in parsed}
        missing = [k for k in field_list if k not in parsed]
        return {
            "ok": True,
            "selected": selected,
            "missing": missing,
        }

    _tool.__name__ = "builtin_json_select_fields"
    return _tool


def _make_files_read_text_tool(project_root: Path, max_chars: int):
    def _tool(path: str, encoding: str = "utf-8") -> dict[str, Any]:
        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        if not resolved.is_file():
            return {"ok": False, "error": f"File not found: {path}"}

        try:
            content = resolved.read_text(encoding=encoding)
            return {
                "ok": True,
                "content": content[:max_chars],
                "truncated": len(content) > max_chars,
                "path": path,
                "size": len(content),
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    _tool.__name__ = "builtin_files_read_text"
    return _tool


def _make_files_write_text_tool(project_root: Path):
    def _tool(path: str, content: str, encoding: str = "utf-8") -> dict[str, Any]:
        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            resolved.parent.mkdir(parents=True, exist_ok=True)
            resolved.write_text(content, encoding=encoding)
            return {
                "ok": True,
                "path": path,
                "size": len(content),
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    _tool.__name__ = "builtin_files_write_text"
    return _tool


def _make_files_list_dir_tool(project_root: Path):
    def _tool(path: str = ".") -> dict[str, Any]:
        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        if not resolved.is_dir():
            return {"ok": False, "error": f"Not a directory: {path}"}

        try:
            entries = []
            for item in sorted(resolved.iterdir()):
                entries.append({
                    "name": item.name,
                    "type": "dir" if item.is_dir() else "file",
                    "size": item.stat().st_size if item.is_file() else 0,
                })
            return {
                "ok": True,
                "path": path,
                "entries": entries[:200],
                "truncated": len(entries) > 200,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    _tool.__name__ = "builtin_files_list_dir"
    return _tool


# ---------------------------------------------------------------------------
# Document processing tools (PDF, PPTX, CSV, DOCX, Markdown, Chart)
# ---------------------------------------------------------------------------

def _create_document_tools(options: DynamicRunOptions) -> list[FunctionTool]:
    """Create document processing tools for reading/writing various formats."""
    project_root = Path(options.project_root or Path.cwd()).expanduser().resolve()
    max_output = max(500, options.shell_output_max_chars or _DEFAULT_MAX_OUTPUT_CHARS)

    return [
        FunctionTool(
            func=_make_pdf_read_text_tool(project_root, max_output),
            name="pdf_read_text",
            description=(
                "Extract text from a PDF file. Returns the text content of "
                "each page. Requires the 'pypdf' package (use pip_install "
                "to install if missing). Path must be relative to the project root."
            ),
        ),
        FunctionTool(
            func=_make_pptx_create_tool(project_root),
            name="pptx_create",
            description=(
                "Create editable .pptx from JSON slide objects. Layouts: "
                "cover, content, comparison, table, cards, metric, timeline, "
                "chart, quote, blank. Supports svg/svg_path; set "
                "engine='ppt-master' for vendored SVG-to-DrawingML native "
                "shape export. Themes: default, dark, editorial, consulting, "
                "academic, tech. Requires python-pptx. Path is project-relative."
            ),
        ),
        FunctionTool(
            func=_make_csv_read_tool(project_root, max_output),
            name="csv_read",
            description=(
                "Read a CSV file and return its rows as a list of dicts "
                "(using the header row as keys). Path must be relative to "
                "the project root. No external dependencies required."
            ),
        ),
        FunctionTool(
            func=_make_csv_write_tool(project_root),
            name="csv_write",
            description=(
                "Write data to a CSV file. Accepts a list of dicts (each "
                "dict is a row, keys become column headers). Path must be "
                "relative to the project root. No external dependencies required."
            ),
        ),
        FunctionTool(
            func=_make_docx_create_tool(project_root),
            name="docx_create",
            description=(
                "Create a Word (.docx) document from structured content. "
                "Accepts a list of content blocks, each with 'type' "
                "('heading', 'paragraph', 'bullets') and 'text' or 'items'. "
                "Requires the 'python-docx' package (use pip_install to "
                "install if missing). Path must be relative to the project root."
            ),
        ),
        FunctionTool(
            func=_make_markdown_write_tool(project_root),
            name="markdown_write",
            description=(
                "Write a Markdown (.md) file. Accepts the markdown content "
                "as a string. Path must be relative to the project root. "
                "No external dependencies required."
            ),
        ),
        FunctionTool(
            func=_make_chart_create_tool(project_root),
            name="chart_create",
            description=(
                "Create a chart image (PNG) from data. Supports chart types: "
                "'bar', 'line', 'pie', 'scatter'. Accepts labels, values, "
                "title, xlabel, ylabel. Requires the 'matplotlib' package "
                "(use pip_install to install if missing). Path must be "
                "relative to the project root."
            ),
        ),
        FunctionTool(
            func=_make_image_create_tool(project_root),
            name="image_create",
            description=(
                "Create an image file. Two modes:\n"
                "1) Programmatic: draw text, rectangles, ellipses, lines on "
                "a canvas using the 'elements' JSON parameter. Requires "
                "'Pillow' package.\n"
                "2) AI-generated: set 'ai_prompt' to generate an image using "
                "OpenAI DALL-E 3 (requires OPENAI_API_KEY env var).\n"
                "Supports .png, .jpg, .webp output. Path must be relative "
                "to the project root."
            ),
        ),
    ]


def _make_pdf_read_text_tool(project_root: Path, max_chars: int):
    def _tool(path: str, pages: str = "") -> dict[str, Any]:
        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        if not resolved.is_file():
            return {"ok": False, "error": f"File not found: {path}"}

        try:
            import pypdf
        except ImportError:
            return {
                "ok": False,
                "error": (
                    "The 'pypdf' package is not installed. "
                    "Use the pip_install tool to install it first: "
                    "pip_install(packages='pypdf')"
                ),
            }

        try:
            reader = pypdf.PdfReader(str(resolved))
            total_pages = len(reader.pages)

            if pages:
                page_indices = _parse_page_range(pages, total_pages)
            else:
                page_indices = list(range(total_pages))

            extracted: list[dict[str, Any]] = []
            total_len = 0
            for i in page_indices:
                text = reader.pages[i].extract_text() or ""
                total_len += len(text)
                extracted.append({"page": i + 1, "text": text})
                if total_len > max_chars:
                    break

            full_text = "\n".join(p["text"] for p in extracted)
            return {
                "ok": True,
                "path": path,
                "total_pages": total_pages,
                "pages_read": len(extracted),
                "text": full_text[:max_chars],
                "truncated": len(full_text) > max_chars,
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to read PDF: {e}"}

    _tool.__name__ = "builtin_pdf_read_text"
    return _tool


def _parse_page_range(pages_str: str, total: int) -> list[int]:
    """Parse page range string like '1-3,5,7-9' into 0-based indices."""
    indices: list[int] = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            s = max(1, int(start.strip()))
            e = min(total, int(end.strip()))
            indices.extend(range(s - 1, e))
        else:
            idx = int(part) - 1
            if 0 <= idx < total:
                indices.append(idx)
    return indices


def _make_pptx_create_tool(project_root: Path):
    def _tool(
        path: str,
        slides: str,
        theme: str = "default",
        engine: str = "auto",
    ) -> dict[str, Any]:
        """Create an editable PPTX deck from structured slide JSON.

        ``engine="ppt-master"`` uses the vendored PPT Master SVG-to-DrawingML
        converter, so slide SVG becomes directly editable native PowerPoint
        shapes. ``engine="python-pptx"`` uses FlowForge's structured fallback
        renderer for native text boxes, shapes, tables, and charts. ``auto``
        chooses PPT Master whenever a slide supplies ``svg`` or ``svg_path``.
        """
        import json as _json

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            slide_data = _json.loads(slides) if isinstance(slides, str) else slides
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid slides JSON: {e}"}

        if not isinstance(slide_data, list):
            return {"ok": False, "error": "slides must be a JSON array of slide objects"}

        engine_name = str(engine or "auto").lower()
        wants_ppt_master = (
            engine_name in {"ppt-master", "ppt_master", "svg", "drawingml"}
            or (
                engine_name == "auto"
                and any(
                    isinstance(slide, dict) and (slide.get("svg") or slide.get("svg_path"))
                    for slide in slide_data
                )
            )
        )
        if wants_ppt_master:
            try:
                from flowforge.tools.ppt_master_bridge import create_native_pptx_from_slide_data
            except Exception as e:
                return {
                    "ok": False,
                    "error": f"PPT Master engine unavailable: {e}",
                }

            try:
                for idx, slide_info in enumerate(slide_data):
                    if not isinstance(slide_info, dict):
                        return {"ok": False, "error": f"slide {idx + 1} must be an object"}
                result = create_native_pptx_from_slide_data(
                    project_root=project_root,
                    output_path=resolved,
                    slide_data=slide_data,
                    theme=theme,
                )
                result.update({
                    "path": path,
                    "slide_count": len(slide_data),
                    "size": resolved.stat().st_size if resolved.exists() else 0,
                })
                if not result.get("ok"):
                    result.setdefault("error", "PPT Master conversion failed")
                return result
            except Exception as e:
                return {"ok": False, "error": f"Failed to create PPTX via PPT Master: {e}"}

        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.dml.color import RGBColor
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError:
            return {
                "ok": False,
                "error": (
                    "The 'python-pptx' package is not installed. "
                    "Use the pip_install tool to install it first: "
                    "pip_install(packages='python-pptx')"
                ),
            }

        palettes = {
            "default": {
                "bg": "F7F8FB", "fg": "172033", "muted": "5F6B7A",
                "accent": "2F6FED", "accent2": "11A683", "panel": "FFFFFF",
                "line": "D8DEE8",
            },
            "dark": {
                "bg": "1E1E2E", "fg": "FFFFFF", "muted": "C9CEDA",
                "accent": "8AB4FF", "accent2": "7AE1B8", "panel": "2B2D42",
                "line": "4A4E69",
            },
            "editorial": {
                "bg": "F6F0E8", "fg": "1C1A17", "muted": "6F6257",
                "accent": "C2472D", "accent2": "243B53", "panel": "FFF9F1",
                "line": "DACBBB",
            },
            "consulting": {
                "bg": "FFFFFF", "fg": "111827", "muted": "4B5563",
                "accent": "0B5CAD", "accent2": "E87500", "panel": "F3F6FA",
                "line": "CBD5E1",
            },
            "academic": {
                "bg": "FAFAF8", "fg": "1F2937", "muted": "56616F",
                "accent": "6A1B9A", "accent2": "0F766E", "panel": "FFFFFF",
                "line": "D6D3D1",
            },
            "tech": {
                "bg": "07111F", "fg": "F8FAFC", "muted": "A7B1C2",
                "accent": "38BDF8", "accent2": "A3E635", "panel": "102033",
                "line": "29445F",
            },
        }

        def _hex(value: Any, fallback: str) -> str:
            text = str(value or fallback).strip().lstrip("#")
            if len(text) == 3:
                text = "".join(ch * 2 for ch in text)
            if len(text) != 6:
                text = fallback
            try:
                int(text, 16)
            except ValueError:
                text = fallback
            return text.upper()

        def _rgb(value: Any, fallback: str = "000000"):
            text = _hex(value, fallback)
            return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))

        def _as_dict(value: Any, fallback_key: str = "title") -> dict[str, Any]:
            if isinstance(value, dict):
                return value
            return {fallback_key: str(value)}

        def _palette(slide_info: dict[str, Any]) -> dict[str, str]:
            base = dict(palettes.get(str(theme).lower(), palettes["default"]))
            slide_theme = str(slide_info.get("theme", "")).lower()
            if slide_theme in palettes:
                base.update(palettes[slide_theme])
            for key in ("bg", "fg", "muted", "accent", "accent2", "panel", "line"):
                if key in slide_info:
                    base[key] = _hex(slide_info[key], base[key])
            return base

        def _emu(value: Any, default: float):
            try:
                return Inches(float(value))
            except (TypeError, ValueError):
                return Inches(default)

        def _set_fill(shape: Any, color: str, transparency: int | None = None) -> None:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(color)
            if transparency is not None:
                fill.transparency = max(0, min(100, int(transparency)))

        def _set_line(shape: Any, color: str | None, width_pt: float = 1.0) -> None:
            if not color:
                shape.line.fill.background()
                return
            shape.line.color.rgb = _rgb(color)
            shape.line.width = Pt(width_pt)

        def _apply_background(slide: Any, pal: dict[str, str]) -> None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(pal["bg"])

        def _set_text_frame(
            shape: Any,
            text: Any,
            *,
            color: str,
            size: float = 20,
            bold: bool = False,
            align: Any = PP_ALIGN.LEFT,
            font: str = "Aptos",
            vertical: Any = MSO_ANCHOR.TOP,
            line_spacing: float | None = None,
        ) -> None:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = vertical
            p = tf.paragraphs[0]
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            run = p.add_run()
            run.text = str(text or "")
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)

        def _text_box(
            slide: Any,
            text: Any,
            x: float,
            y: float,
            w: float,
            h: float,
            *,
            color: str,
            size: float = 20,
            bold: bool = False,
            align: Any = PP_ALIGN.LEFT,
            font: str = "Aptos",
            vertical: Any = MSO_ANCHOR.TOP,
        ) -> Any:
            shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            _set_text_frame(
                shape, text, color=color, size=size, bold=bold,
                align=align, font=font, vertical=vertical,
            )
            return shape

        def _add_title(slide: Any, info: dict[str, Any], pal: dict[str, str]) -> None:
            title = info.get("title", "")
            if title:
                _text_box(
                    slide, title, 0.65, 0.35, 8.6, 0.62,
                    color=pal["fg"], size=25, bold=True,
                )
            kicker = info.get("kicker") or info.get("eyebrow")
            if kicker:
                _text_box(
                    slide, str(kicker).upper(), 0.68, 0.16, 5.6, 0.24,
                    color=pal["accent"], size=8.5, bold=True,
                )

        def _add_footer(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            footer = info.get("footer")
            if footer is None:
                footer = ""
            if footer or info.get("show_page_number"):
                text = str(footer)
                if info.get("show_page_number"):
                    text = f"{text}  {idx + 1}".strip()
                _text_box(
                    slide, text, 0.65, 7.05, 12.0, 0.22,
                    color=pal["muted"], size=7.5,
                )

        def _bullets(
            slide: Any,
            bullets: list[Any],
            x: float,
            y: float,
            w: float,
            h: float,
            pal: dict[str, str],
            *,
            size: float = 18,
        ) -> None:
            shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(bullet)
                p.level = 0
                p.space_after = Pt(8)
                p.font.name = "Aptos"
                p.font.size = Pt(size)
                p.font.color.rgb = _rgb(pal["fg"])

        def _shape_rect(
            slide: Any,
            x: float,
            y: float,
            w: float,
            h: float,
            *,
            fill: str,
            line: str | None = None,
            radius: bool = False,
        ) -> Any:
            kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
            shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
            _set_fill(shape, fill)
            _set_line(shape, line, 0.8)
            return shape

        def _add_image(slide: Any, info: dict[str, Any]) -> None:
            image_path = info.get("image_path", "")
            if not image_path:
                return
            try:
                img_resolved = _resolve_safe_path(project_root, image_path)
            except Exception:
                return
            if not img_resolved.is_file():
                return
            x = float(info.get("image_x", 7.1))
            y = float(info.get("image_y", 1.35))
            w = float(info.get("image_w", 5.35))
            h = info.get("image_h")
            if h is None:
                slide.shapes.add_picture(str(img_resolved), Inches(x), Inches(y), width=Inches(w))
            else:
                slide.shapes.add_picture(
                    str(img_resolved), Inches(x), Inches(y),
                    width=Inches(w), height=Inches(float(h)),
                )

        def _render_cover(slide: Any, info: dict[str, Any], pal: dict[str, str]) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _shape_rect(slide, 0, 0, 0.18, 7.5, fill=pal["accent"], line=None)
            _shape_rect(slide, 9.4, 0, 3.95, 7.5, fill=pal["panel"], line=None)
            if info.get("kicker") or info.get("eyebrow"):
                _text_box(
                    slide, str(info.get("kicker") or info.get("eyebrow")).upper(),
                    0.75, 1.05, 5.8, 0.28, color=pal["accent"], size=9, bold=True,
                )
            _text_box(
                slide, info.get("title", ""), 0.72, 1.55, 8.15, 1.9,
                color=pal["fg"], size=34, bold=True,
            )
            if info.get("subtitle") or info.get("body"):
                _text_box(
                    slide, info.get("subtitle") or info.get("body"), 0.78, 3.65, 7.3, 0.8,
                    color=pal["muted"], size=16,
                )
            if info.get("bullets"):
                _bullets(slide, list(info["bullets"])[:4], 0.9, 4.8, 6.8, 1.25, pal, size=13)
            _add_image(slide, {**info, "image_x": info.get("image_x", 9.85), "image_y": info.get("image_y", 1.35), "image_w": info.get("image_w", 2.85)})

        def _render_section(slide: Any, info: dict[str, Any], pal: dict[str, str]) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _shape_rect(slide, 0.68, 1.45, 0.1, 4.2, fill=pal["accent"], line=None)
            _text_box(slide, info.get("title", ""), 1.05, 2.35, 9.6, 1.25, color=pal["fg"], size=32, bold=True)
            if info.get("subtitle") or info.get("body"):
                _text_box(slide, info.get("subtitle") or info.get("body"), 1.08, 3.72, 8.8, 0.8, color=pal["muted"], size=16)

        def _render_content(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            if info.get("bullets"):
                _bullets(slide, list(info["bullets"]), 0.88, 1.45, 6.4, 4.8, pal, size=16)
            elif info.get("body"):
                _text_box(slide, info["body"], 0.78, 1.35, 6.1, 4.85, color=pal["fg"], size=16)
            _add_image(slide, info)
            _render_shapes(slide, info, pal)
            _add_footer(slide, info, pal, idx)

        def _render_comparison(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            columns = [info.get("left", {}), info.get("right", {})]
            for i, col in enumerate(columns):
                col = _as_dict(col, "heading")
                x = 0.78 + i * 6.18
                _shape_rect(slide, x, 1.35, 5.62, 5.35, fill=pal["panel"], line=pal["line"], radius=True)
                _text_box(slide, col.get("heading", f"Option {i + 1}"), x + 0.35, 1.72, 4.8, 0.35, color=pal["accent" if i == 0 else "accent2"], size=17, bold=True)
                _bullets(slide, list(col.get("bullets", [])), x + 0.42, 2.35, 4.72, 3.6, pal, size=13)
            _add_footer(slide, info, pal, idx)

        def _render_table(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            tbl_data = info.get("table", {}) or {}
            headers = list(tbl_data.get("headers", []))
            rows = list(tbl_data.get("rows", []))
            if not headers:
                return
            n_rows = len(rows) + 1
            n_cols = len(headers)
            x = float(tbl_data.get("x", 0.7))
            y = float(tbl_data.get("y", 1.35))
            w = float(tbl_data.get("w", 11.95))
            h = float(tbl_data.get("h", min(5.7, 0.46 * n_rows + 0.2)))
            table = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(pal["accent"])
                for para in cell.text_frame.paragraphs:
                    para.font.bold = True
                    para.font.size = Pt(10)
                    para.font.color.rgb = _rgb("FFFFFF")
            for row_idx, row in enumerate(rows):
                for col_idx in range(n_cols):
                    value = row[col_idx] if isinstance(row, (list, tuple)) and col_idx < len(row) else ""
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(value)
                    if row_idx % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(pal["panel"])
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(9.2)
                        para.font.color.rgb = _rgb(pal["fg"])
            _add_footer(slide, info, pal, idx)

        def _render_cards(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            cards = list(info.get("cards") or [])
            if not cards and info.get("bullets"):
                cards = [{"title": str(b), "body": ""} for b in info["bullets"]]
            n = max(1, min(len(cards), 4))
            card_w = 11.9 / n
            for i, raw_card in enumerate(cards[:4]):
                card = _as_dict(raw_card)
                x = 0.72 + i * card_w
                _shape_rect(slide, x, 1.55, card_w - 0.22, 4.75, fill=pal["panel"], line=pal["line"], radius=True)
                _text_box(slide, card.get("title", ""), x + 0.25, 1.88, card_w - 0.72, 0.72, color=pal["accent"], size=15, bold=True)
                if card.get("body"):
                    _text_box(slide, card.get("body", ""), x + 0.25, 2.72, card_w - 0.72, 1.4, color=pal["fg"], size=11.5)
                if card.get("bullets"):
                    _bullets(slide, list(card["bullets"])[:4], x + 0.28, 4.15, card_w - 0.75, 1.55, pal, size=10.5)
            _add_footer(slide, info, pal, idx)

        def _render_metric(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            metrics = list(info.get("metrics") or [])
            if not metrics and info.get("bullets"):
                metrics = [{"value": str(b), "label": ""} for b in info["bullets"][:3]]
            n = max(1, min(len(metrics), 3))
            for i, raw_metric in enumerate(metrics[:3]):
                metric = _as_dict(raw_metric, "value")
                x = 0.85 + i * 4.05
                _text_box(slide, metric.get("value", ""), x, 2.05, 3.55, 0.85, color=pal["accent" if i == 0 else "fg"], size=30, bold=True)
                _text_box(slide, metric.get("label", ""), x + 0.03, 3.02, 3.25, 0.42, color=pal["fg"], size=13, bold=True)
                note = metric.get("note") or metric.get("delta") or ""
                if note:
                    _text_box(slide, note, x + 0.03, 3.55, 3.35, 0.65, color=pal["muted"], size=10.5)
            if info.get("body"):
                _text_box(slide, info["body"], 0.9, 5.12, 10.8, 0.8, color=pal["muted"], size=14)
            _add_footer(slide, info, pal, idx)

        def _render_quote(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _text_box(slide, "“", 0.78, 0.95, 0.8, 0.8, color=pal["accent"], size=42, bold=True)
            _text_box(slide, info.get("quote") or info.get("body") or info.get("title", ""), 1.25, 1.55, 10.4, 2.35, color=pal["fg"], size=27, bold=True)
            attribution = info.get("attribution") or info.get("subtitle", "")
            if attribution:
                _text_box(slide, attribution, 1.33, 4.22, 7.8, 0.42, color=pal["muted"], size=13)
            _add_footer(slide, info, pal, idx)

        def _render_timeline(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            items = list(info.get("items") or info.get("timeline") or info.get("steps") or [])
            if not items and info.get("bullets"):
                items = [{"title": str(b)} for b in info["bullets"]]
            y = 3.1
            _shape_rect(slide, 1.05, y + 0.12, 11.1, 0.03, fill=pal["line"], line=None)
            n = max(1, min(len(items), 5))
            for i, raw_item in enumerate(items[:5]):
                item = _as_dict(raw_item)
                x = 1.05 + i * (11.1 / max(1, n - 1))
                oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(y), Inches(0.24), Inches(0.24))
                _set_fill(oval, pal["accent"])
                _set_line(oval, None)
                _text_box(slide, item.get("label") or item.get("date") or str(i + 1), x - 0.45, y - 0.5, 0.9, 0.3, color=pal["accent"], size=9, bold=True, align=PP_ALIGN.CENTER)
                _text_box(slide, item.get("title", ""), x - 0.92, y + 0.42, 1.85, 0.55, color=pal["fg"], size=11, bold=True, align=PP_ALIGN.CENTER)
                if item.get("body"):
                    _text_box(slide, item["body"], x - 0.95, y + 1.05, 1.9, 0.72, color=pal["muted"], size=8.8, align=PP_ALIGN.CENTER)
            _add_footer(slide, info, pal, idx)

        def _render_chart(slide: Any, info: dict[str, Any], pal: dict[str, str], idx: int) -> None:
            _shape_rect(slide, 0, 0, 13.333, 7.5, fill=pal["bg"], line=None)
            _add_title(slide, info, pal)
            chart_info = info.get("chart", {}) or {}
            categories = list(chart_info.get("categories") or [])
            series = chart_info.get("series")
            values = chart_info.get("values")
            if not series and values is not None:
                series = [{"name": chart_info.get("name", "Series 1"), "values": values}]
            elif isinstance(series, dict):
                series = [series]
            if not categories or not series:
                if info.get("bullets"):
                    _bullets(slide, list(info["bullets"]), 0.9, 1.45, 10.5, 4.7, pal)
                return
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for raw_serie in series:
                serie = _as_dict(raw_serie, "name")
                chart_data.add_series(str(serie.get("name", "Series")), list(serie.get("values", [])))
            chart_types = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
            }
            chart_type = chart_types.get(str(chart_info.get("type", "column")).lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            frame = slide.shapes.add_chart(
                chart_type,
                _emu(chart_info.get("x"), 0.85),
                _emu(chart_info.get("y"), 1.45),
                _emu(chart_info.get("w"), 11.8),
                _emu(chart_info.get("h"), 4.9),
                chart_data,
            )
            chart = frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.has_title = bool(chart_info.get("title"))
            if chart_info.get("title"):
                chart.chart_title.text_frame.text = str(chart_info["title"])
            _add_footer(slide, info, pal, idx)

        def _render_shapes(slide: Any, info: dict[str, Any], pal: dict[str, str]) -> None:
            shape_map = {
                "rect": MSO_SHAPE.RECTANGLE,
                "rectangle": MSO_SHAPE.RECTANGLE,
                "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_SHAPE.OVAL,
                "oval": MSO_SHAPE.OVAL,
            }
            for obj in list(info.get("shapes") or info.get("objects") or []):
                kind = shape_map.get(str(obj.get("type", "rect")).lower(), MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(
                    kind,
                    _emu(obj.get("x"), 0),
                    _emu(obj.get("y"), 0),
                    _emu(obj.get("w"), 1),
                    _emu(obj.get("h"), 1),
                )
                _set_fill(shape, obj.get("fill", pal["panel"]), obj.get("transparency"))
                _set_line(shape, obj.get("stroke", obj.get("line", pal["line"])), float(obj.get("stroke_width", 1)))
                if obj.get("text"):
                    _set_text_frame(
                        shape, obj["text"],
                        color=_hex(obj.get("color"), pal["fg"]),
                        size=float(obj.get("font_size", 14)),
                        bold=bool(obj.get("bold", False)),
                        align=PP_ALIGN.CENTER,
                        vertical=MSO_ANCHOR.MIDDLE,
                    )

        renderers = {
            "cover": _render_cover,
            "section": _render_section,
            "content": _render_content,
            "comparison": _render_comparison,
            "table": _render_table,
            "cards": _render_cards,
            "metric": _render_metric,
            "quote": _render_quote,
            "timeline": _render_timeline,
            "process": _render_timeline,
            "chart": _render_chart,
            "blank": _render_content,
        }

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]

            for idx, slide_info in enumerate(slide_data):
                if not isinstance(slide_info, dict):
                    return {"ok": False, "error": f"slide {idx + 1} must be an object"}
                slide = prs.slides.add_slide(blank_layout)
                pal = _palette(slide_info)
                _apply_background(slide, pal)
                layout_name = str(slide_info.get("layout", "content")).lower()
                renderer = renderers.get(layout_name, _render_content)
                if renderer in {_render_cover, _render_section}:
                    renderer(slide, slide_info, pal)
                else:
                    renderer(slide, slide_info, pal, idx)

                speaker_note = slide_info.get("speaker_note", slide_info.get("notes", ""))
                if speaker_note:
                    slide.notes_slide.notes_text_frame.text = str(speaker_note)

            resolved.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(resolved))
            return {
                "ok": True,
                "path": path,
                "slide_count": len(slide_data),
                "size": resolved.stat().st_size,
                "engine": "python-pptx",
                "native_objects": True,
                "layouts": [
                    str(slide.get("layout", "content")).lower()
                    for slide in slide_data
                    if isinstance(slide, dict)
                ],
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to create PPTX: {e}"}

    _tool.__name__ = "builtin_pptx_create"
    return _tool


def _make_csv_read_tool(project_root: Path, max_chars: int):
    def _tool(path: str, encoding: str = "utf-8", max_rows: int = 500) -> dict[str, Any]:
        import csv as _csv

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        if not resolved.is_file():
            return {"ok": False, "error": f"File not found: {path}"}

        try:
            with resolved.open(encoding=encoding, newline="") as f:
                reader = _csv.DictReader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i >= max_rows:
                        break
                    rows.append(dict(row))
                headers = reader.fieldnames or []
            return {
                "ok": True,
                "path": path,
                "headers": list(headers),
                "row_count": len(rows),
                "rows": rows,
                "truncated": len(rows) >= max_rows,
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to read CSV: {e}"}

    _tool.__name__ = "builtin_csv_read"
    return _tool


def _make_csv_write_tool(project_root: Path):
    def _tool(path: str, data: str, encoding: str = "utf-8") -> dict[str, Any]:
        import csv as _csv
        import json as _json

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            rows = _json.loads(data) if isinstance(data, str) else data
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid data JSON: {e}"}

        if not isinstance(rows, list) or not rows:
            return {"ok": False, "error": "data must be a non-empty JSON array of objects"}

        try:
            resolved.parent.mkdir(parents=True, exist_ok=True)
            headers = list(rows[0].keys())
            with resolved.open("w", encoding=encoding, newline="") as f:
                writer = _csv.DictWriter(f, fieldnames=headers)
                writer.writeheader()
                writer.writerows(rows)
            return {
                "ok": True,
                "path": path,
                "row_count": len(rows),
                "headers": headers,
                "size": resolved.stat().st_size,
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to write CSV: {e}"}

    _tool.__name__ = "builtin_csv_write"
    return _tool


def _make_docx_create_tool(project_root: Path):
    def _tool(path: str, content: str) -> dict[str, Any]:
        import json as _json

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            blocks = _json.loads(content) if isinstance(content, str) else content
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid content JSON: {e}"}

        if not isinstance(blocks, list):
            return {"ok": False, "error": "content must be a JSON array of block objects"}

        try:
            from docx import Document
        except ImportError:
            return {
                "ok": False,
                "error": (
                    "The 'python-docx' package is not installed. "
                    "Use the pip_install tool to install it first: "
                    "pip_install(packages='python-docx')"
                ),
            }

        try:
            doc = Document()
            for block in blocks:
                block_type = block.get("type", "paragraph")
                if block_type == "heading":
                    level = block.get("level", 1)
                    doc.add_heading(block.get("text", ""), level=level)
                elif block_type == "paragraph":
                    doc.add_paragraph(block.get("text", ""))
                elif block_type == "bullets":
                    for item in block.get("items", []):
                        doc.add_paragraph(str(item), style="List Bullet")
                elif block_type == "table":
                    rows_data = block.get("rows", [])
                    headers = block.get("headers", [])
                    if headers:
                        table = doc.add_table(rows=1, cols=len(headers))
                        table.style = "Table Grid"
                        for i, h in enumerate(headers):
                            table.rows[0].cells[i].text = str(h)
                        for row_data in rows_data:
                            row = table.add_row()
                            for i, val in enumerate(row_data):
                                if i < len(headers):
                                    row.cells[i].text = str(val)

            resolved.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(resolved))
            return {
                "ok": True,
                "path": path,
                "block_count": len(blocks),
                "size": resolved.stat().st_size,
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to create DOCX: {e}"}

    _tool.__name__ = "builtin_docx_create"
    return _tool


def _make_markdown_write_tool(project_root: Path):
    def _tool(path: str, content: str) -> dict[str, Any]:
        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            resolved.parent.mkdir(parents=True, exist_ok=True)
            resolved.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": path,
                "size": len(content),
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to write Markdown: {e}"}

    _tool.__name__ = "builtin_markdown_write"
    return _tool


def _make_chart_create_tool(project_root: Path):
    def _tool(
        path: str,
        chart_type: str,
        labels: str,
        values: str,
        title: str = "",
        xlabel: str = "",
        ylabel: str = "",
    ) -> dict[str, Any]:
        import json as _json

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        try:
            label_list = _json.loads(labels) if isinstance(labels, str) else labels
            value_list = _json.loads(values) if isinstance(values, str) else values
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid JSON for labels/values: {e}"}

        if chart_type not in ("bar", "line", "pie", "scatter"):
            return {
                "ok": False,
                "error": f"Unsupported chart_type: {chart_type!r}. "
                         f"Use 'bar', 'line', 'pie', or 'scatter'.",
            }

        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return {
                "ok": False,
                "error": (
                    "The 'matplotlib' package is not installed. "
                    "Use the pip_install tool to install it first: "
                    "pip_install(packages='matplotlib')"
                ),
            }

        try:
            fig, ax = plt.subplots(figsize=(10, 6))
            if chart_type == "bar":
                ax.bar(label_list, value_list)
            elif chart_type == "line":
                ax.plot(label_list, value_list, marker="o")
            elif chart_type == "pie":
                ax.pie(value_list, labels=label_list, autopct="%1.1f%%")
            elif chart_type == "scatter":
                ax.scatter(label_list, value_list)

            if title:
                ax.set_title(title)
            if xlabel:
                ax.set_xlabel(xlabel)
            if ylabel:
                ax.set_ylabel(ylabel)

            plt.tight_layout()
            resolved.parent.mkdir(parents=True, exist_ok=True)
            fig.savefig(str(resolved), dpi=150)
            plt.close(fig)

            return {
                "ok": True,
                "path": path,
                "chart_type": chart_type,
                "size": resolved.stat().st_size,
            }
        except Exception as e:
            plt.close("all")
            return {"ok": False, "error": f"Failed to create chart: {e}"}

    _tool.__name__ = "builtin_chart_create"
    return _tool


def _make_image_create_tool(project_root: Path):
    """Create images programmatically using Pillow.

    Supports: blank canvas, text rendering, basic shapes, compositing.
    For AI-generated images, uses OpenAI DALL-E when ``OPENAI_API_KEY`` is set.
    """

    def _tool(
        path: str,
        width: int = 1024,
        height: int = 768,
        background: str = "#FFFFFF",
        elements: str = "[]",
        ai_prompt: str = "",
        ai_provider: str = "openai",
        ai_size: str = "1024x1024",
    ) -> dict[str, Any]:
        """Create an image file.

        Parameters
        ----------
        path : str
            Output file path (relative to project root). Supports .png, .jpg, .webp.
        width / height : int
            Canvas dimensions (ignored when ``ai_prompt`` is set).
        background : str
            Background colour (hex like ``"#FF0000"`` or name like ``"white"``).
        elements : str
            JSON array of drawing elements. Each element is an object with:
            - ``"type"``: ``"text"``, ``"rectangle"``, ``"ellipse"``, ``"line"``
            - ``"text"`` type: ``text``, ``x``, ``y``, ``size`` (font size),
              ``color``, ``bold`` (bool)
            - ``"rectangle"`` type: ``x``, ``y``, ``w``, ``h``, ``color``,
              ``fill`` (optional fill colour)
            - ``"ellipse"`` type: ``x``, ``y``, ``w``, ``h``, ``color``, ``fill``
            - ``"line"`` type: ``x1``, ``y1``, ``x2``, ``y2``, ``color``,
              ``width`` (line width)
        ai_prompt : str
            When non-empty, generate the image using an AI model instead of
            drawing elements.  Requires ``OPENAI_API_KEY`` env var for OpenAI
            DALL-E.
        ai_provider : str
            AI provider: ``"openai"`` (default).
        ai_size : str
            AI image size: ``"1024x1024"``, ``"1792x1024"``, ``"1024x1792"``.
        """
        import json as _json

        try:
            resolved = _resolve_safe_path(project_root, path)
        except ValueError as e:
            return {"ok": False, "error": str(e)}

        # ── AI image generation ──────────────────────────────────────
        if ai_prompt:
            return _generate_ai_image(
                resolved, ai_prompt, ai_provider, ai_size,
            )

        # ── Programmatic image creation with Pillow ──────────────────
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            return {
                "ok": False,
                "error": (
                    "The 'Pillow' package is not installed. "
                    "Use the pip_install tool to install it first: "
                    "pip_install(packages='Pillow')"
                ),
            }

        try:
            element_list = _json.loads(elements) if isinstance(elements, str) else elements
        except (_json.JSONDecodeError, TypeError) as e:
            return {"ok": False, "error": f"Invalid elements JSON: {e}"}

        if not isinstance(element_list, list):
            element_list = []

        try:
            img = Image.new("RGBA", (width, height), background)
            draw = ImageDraw.Draw(img)

            for elem in element_list:
                etype = elem.get("type", "")

                if etype == "text":
                    text = elem.get("text", "")
                    x, y = elem.get("x", 0), elem.get("y", 0)
                    color = elem.get("color", "#000000")
                    font_size = elem.get("size", 24)
                    try:
                        font = ImageFont.truetype("arial.ttf", font_size)
                    except (OSError, IOError):
                        try:
                            font = ImageFont.truetype(
                                "/System/Library/Fonts/Helvetica.ttc",
                                font_size,
                            )
                        except (OSError, IOError):
                            font = ImageFont.load_default()
                    draw.text((x, y), text, fill=color, font=font)

                elif etype == "rectangle":
                    x, y = elem.get("x", 0), elem.get("y", 0)
                    w, h = elem.get("w", 100), elem.get("h", 100)
                    color = elem.get("color", "#000000")
                    fill = elem.get("fill", None)
                    draw.rectangle([x, y, x + w, y + h], outline=color, fill=fill)

                elif etype == "ellipse":
                    x, y = elem.get("x", 0), elem.get("y", 0)
                    w, h = elem.get("w", 100), elem.get("h", 100)
                    color = elem.get("color", "#000000")
                    fill = elem.get("fill", None)
                    draw.ellipse([x, y, x + w, y + h], outline=color, fill=fill)

                elif etype == "line":
                    x1, y1 = elem.get("x1", 0), elem.get("y1", 0)
                    x2, y2 = elem.get("x2", 100), elem.get("y2", 100)
                    color = elem.get("color", "#000000")
                    lw = elem.get("width", 2)
                    draw.line([(x1, y1), (x2, y2)], fill=color, width=lw)

            # Save
            resolved.parent.mkdir(parents=True, exist_ok=True)
            # Convert RGBA to RGB for JPEG
            ext = resolved.suffix.lower()
            if ext in (".jpg", ".jpeg"):
                img = img.convert("RGB")
            img.save(str(resolved))

            return {
                "ok": True,
                "path": path,
                "width": width,
                "height": height,
                "element_count": len(element_list),
                "size": resolved.stat().st_size,
            }
        except Exception as e:
            return {"ok": False, "error": f"Failed to create image: {e}"}

    _tool.__name__ = "builtin_image_create"
    return _tool


def _generate_ai_image(
    resolved: Path,
    prompt: str,
    provider: str,
    size: str,
) -> dict[str, Any]:
    """Generate an image using an AI provider API."""
    import os as _os

    if provider == "openai":
        api_key = _os.environ.get("OPENAI_API_KEY", "")
        if not api_key:
            return {
                "ok": False,
                "error": (
                    "OPENAI_API_KEY environment variable is not set. "
                    "Set it to use AI image generation with OpenAI DALL-E."
                ),
            }

        try:
            import urllib.request
            import urllib.error
            import json as _json
            import base64

            req_body = _json.dumps({
                "model": "dall-e-3",
                "prompt": prompt,
                "n": 1,
                "size": size,
                "response_format": "b64_json",
            }).encode("utf-8")

            req = urllib.request.Request(
                "https://api.openai.com/v1/images/generations",
                data=req_body,
                headers={
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {api_key}",
                },
                method="POST",
            )

            with urllib.request.urlopen(req, timeout=120) as resp:
                result = _json.loads(resp.read().decode("utf-8"))

            b64_data = result["data"][0]["b64_json"]
            img_bytes = base64.b64decode(b64_data)

            resolved.parent.mkdir(parents=True, exist_ok=True)
            resolved.write_bytes(img_bytes)

            return {
                "ok": True,
                "path": str(resolved.name),
                "provider": "openai",
                "model": "dall-e-3",
                "size": size,
                "file_size": len(img_bytes),
                "revised_prompt": result["data"][0].get("revised_prompt", ""),
            }
        except urllib.error.HTTPError as e:
            body = e.read().decode("utf-8", errors="replace")
            return {"ok": False, "error": f"OpenAI API error {e.code}: {body[:500]}"}
        except Exception as e:
            return {"ok": False, "error": f"AI image generation failed: {e}"}

    return {"ok": False, "error": f"Unsupported AI image provider: {provider!r}"}


def _resolve_safe_path(project_root: Path, path: str) -> Path:
    """Resolve a relative path, ensuring it stays inside project_root."""
    root = project_root.expanduser().resolve()
    raw = Path(path).expanduser()
    resolved = raw.resolve() if raw.is_absolute() else (root / raw).resolve()
    try:
        resolved.relative_to(root)
    except ValueError:
        raise ValueError(f"path must stay inside project root: {resolved}")
    return resolved


def _make_claude_skill_tool(options: DynamicRunOptions):
    """Create a tool that invokes Claude Code skills via the CLI."""

    def _tool(
        skill_name: str,
        prompt: str,
        timeout_seconds: int = 300,
        model: str = "",
    ) -> dict[str, Any]:
        """Invoke a Claude Code skill and return the result.

        Parameters
        ----------
        skill_name : str
            The skill to invoke (e.g. ``"commit"``, ``"review-pr"``).
        prompt : str
            The prompt/arguments for the skill.
        timeout_seconds : int
            Maximum execution time (default 300s).
        model : str
            Model override (e.g. ``"sonnet"``).
        """
        import shutil

        claude_bin = shutil.which("claude")
        if not claude_bin:
            return {
                "ok": False,
                "error": (
                    "Claude CLI ('claude') not found in PATH. "
                    "Install Claude Code: https://docs.anthropic.com/en/docs/claude-code"
                ),
                "skill": skill_name,
            }

        full_prompt = f"/{skill_name} {prompt}"
        cmd = [claude_bin, "--print", "--no-input"]
        if model:
            cmd.extend(["--model", model])
        cmd.extend(["--max-tokens", "4096"])
        cmd.extend(["--prompt", full_prompt])

        run_cwd = options.project_root or None

        try:
            completed = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=timeout_seconds,
                cwd=run_cwd,
                check=False,
            )

            max_output = max(500, options.shell_output_max_chars or _DEFAULT_MAX_OUTPUT_CHARS)
            if completed.returncode == 0:
                return {
                    "ok": True,
                    "result": _trim_output(completed.stdout.strip(), max_output),
                    "skill": skill_name,
                }
            else:
                return {
                    "ok": False,
                    "error": _trim_output(
                        completed.stderr.strip() or f"exit code {completed.returncode}",
                        max_output,
                    ),
                    "result": _trim_output(completed.stdout.strip(), max_output),
                    "skill": skill_name,
                }
        except subprocess.TimeoutExpired:
            return {
                "ok": False,
                "error": f"Skill '{skill_name}' timed out after {timeout_seconds}s",
                "skill": skill_name,
            }
        except Exception as e:
            return {
                "ok": False,
                "error": f"Failed to invoke skill: {e}",
                "skill": skill_name,
            }

    _tool.__name__ = "builtin_claude_skill"
    return _tool


def _make_import_check_tool():
    def _tool(module_names: str, package_names: str = "") -> dict[str, Any]:
        import importlib.metadata
        import importlib.util

        modules = _split_names(module_names)
        packages = _split_names(package_names)
        results: list[dict[str, Any]] = []

        for index, module_name in enumerate(modules):
            package_name = packages[index] if index < len(packages) else module_name
            spec = importlib.util.find_spec(module_name)
            version = None
            if spec is not None:
                try:
                    version = importlib.metadata.version(package_name)
                except importlib.metadata.PackageNotFoundError:
                    version = None
            results.append({
                "module": module_name,
                "package": package_name,
                "available": spec is not None,
                "version": version,
            })

        return {"ok": all(item["available"] for item in results), "results": results}

    _tool.__name__ = "builtin_python_import_check"
    return _tool


def _make_pip_install_tool(options: DynamicRunOptions):
    def _tool(packages: str) -> dict[str, Any]:
        policy = options.dependency_policy
        if not policy.allow_install:
            return {
                "ok": False,
                "error": (
                    "Package installation is disabled. Set "
                    "DependencyPolicy(allow_install=True) to enable it."
                ),
            }
        if "pip" not in policy.allowed_managers:
            return {
                "ok": False,
                "error": "pip is not in allowed_managers.",
            }

        names = _split_names(packages)
        if not names:
            return {"ok": False, "error": "No package names provided."}

        # Validate against the dependency policy.
        if policy.allowed_packages:
            disallowed = [p for p in names if p not in policy.allowed_packages]
            if disallowed:
                return {
                    "ok": False,
                    "error": f"Packages not in allowed_packages: {disallowed}",
                }
        denied = [p for p in names if p in policy.denied_packages]
        if denied:
            return {
                "ok": False,
                "error": f"Packages denied by policy: {denied}",
            }

        import sys

        cmd = [sys.executable, "-m", "pip", "install", "--quiet", *names]
        try:
            completed = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=options.shell_timeout_seconds,
                check=False,
            )
            max_output = max(500, options.shell_output_max_chars or _DEFAULT_MAX_OUTPUT_CHARS)
            return {
                "ok": completed.returncode == 0,
                "packages": names,
                "returncode": completed.returncode,
                "stdout": _trim_output(completed.stdout, max_output),
                "stderr": _trim_output(completed.stderr, max_output),
            }
        except subprocess.TimeoutExpired:
            return {
                "ok": False,
                "packages": names,
                "error": "pip install timed out",
            }

    _tool.__name__ = "builtin_pip_install"
    return _tool


def _make_mcp_start_tool(options: DynamicRunOptions):
    def _tool(server_name: str, cwd: str | None = None) -> dict[str, Any]:
        return _start_declared_mcp_server(options, server_name, cwd=cwd)

    _tool.__name__ = "builtin_mcp_start_server"
    return _tool


def _mcp_endpoint_ready(url: str, timeout_seconds: float = 0.5) -> bool:
    """Return True when an MCP HTTP endpoint host/port accepts connections."""
    if not url:
        return False
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        return False
    port = parsed.port or (443 if parsed.scheme == "https" else 80)
    try:
        with socket.create_connection(
            (parsed.hostname, port),
            timeout=timeout_seconds,
        ):
            return True
    except OSError:
        return False


def _wait_for_mcp_endpoint(
    url: str,
    proc: subprocess.Popen[Any] | None,
    timeout_seconds: int,
) -> bool:
    """Wait until a local MCP endpoint is reachable or the process exits."""
    deadline = time.monotonic() + max(0, timeout_seconds)
    while time.monotonic() <= deadline:
        if _mcp_endpoint_ready(url):
            return True
        if proc is not None and proc.poll() is not None:
            return False
        time.sleep(0.2)
    return _mcp_endpoint_ready(url)


def _start_declared_mcp_server(
    options: DynamicRunOptions,
    server_name: str,
    cwd: str | None = None,
) -> dict[str, Any]:
    command = options.mcp_server_commands.get(server_name)
    endpoint = options.mcp_server_urls.get(server_name, "")
    if not command:
        return {
            "ok": False,
            "server_name": server_name,
            "url": endpoint,
            "already_running": _mcp_endpoint_ready(endpoint),
            "stderr": "MCP server is not declared in DynamicRunOptions.",
        }

    if endpoint and _mcp_endpoint_ready(endpoint):
        return {
            "ok": True,
            "server_name": server_name,
            "pid": None,
            "command": command,
            "cwd": None,
            "url": endpoint,
            "already_running": True,
            "stderr": "",
        }

    project_root = Path(options.project_root or Path.cwd()).expanduser().resolve()
    run_cwd = _resolve_cwd(project_root, cwd)
    proc = subprocess.Popen(
        command,
        cwd=run_cwd,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        start_new_session=True,
    )
    ready = (
        _wait_for_mcp_endpoint(endpoint, proc, options.mcp_start_timeout_seconds)
        if endpoint
        else proc.poll() is None
    )
    return {
        "ok": ready,
        "server_name": server_name,
        "pid": proc.pid,
        "command": command,
        "cwd": str(run_cwd),
        "url": endpoint,
        "already_running": False,
        "stderr": "" if ready else "MCP server command exited or endpoint was not ready.",
    }


def _make_mcp_register_tool(options: DynamicRunOptions):
    def _tool(
        server_name: str,
        tool_names: str = "",
        url: str = "",
        description: str = "",
        auto_start: bool = True,
        cwd: str | None = None,
        ctx: Any = None,
    ) -> dict[str, Any]:
        if ctx is None:
            return {
                "ok": False,
                "server_name": server_name,
                "stderr": "mcp_register_server must be called via ctx.call_tool().",
            }

        endpoint = url or options.mcp_server_urls.get(server_name, "")
        if not endpoint:
            return {
                "ok": False,
                "server_name": server_name,
                "stderr": (
                    "MCP server URL is not declared. Pass url=... or set "
                    "DynamicRunOptions.mcp_server_urls[server_name]."
                ),
            }

        start_result: dict[str, Any] | None = None
        if (
            auto_start
            and server_name in options.mcp_server_commands
            and not _mcp_endpoint_ready(endpoint)
        ):
            start_result = _start_declared_mcp_server(
                options, server_name, cwd=cwd,
            )
            if not start_result.get("ok"):
                return {
                    "ok": False,
                    "server_name": server_name,
                    "url": endpoint,
                    "start": start_result,
                    "stderr": (
                        "MCP server could not be started before registration."
                    ),
                }

        declared = options.mcp_server_tools.get(server_name, [])
        raw_names = tool_names or ",".join(declared)
        names = [
            item.strip()
            for item in raw_names.replace("\n", ",").split(",")
            if item.strip()
        ]
        if not names:
            return {
                "ok": False,
                "server_name": server_name,
                "url": endpoint,
                "stderr": (
                    "No MCP tool names were provided. Pass tool_names as a "
                    "comma-separated string or set DynamicRunOptions.mcp_server_tools."
                ),
            }

        from flowforge.types import MCPServer

        registered: list[str] = []
        existing = {
            getattr(tool, "name", "")
            for tool in getattr(ctx.global_ctx, "global_tools", [])
        }
        for name in names:
            if name in existing:
                registered.append(name)
                continue
            ctx.global_ctx.global_tools.append(
                MCPServer(
                    url=endpoint,
                    name=name,
                    description=description or f"{server_name} MCP tool: {name}",
                    headers=options.mcp_server_headers.get(server_name, {}),
                )
            )
            existing.add(name)
            registered.append(name)

        return {
            "ok": True,
            "server_name": server_name,
            "url": endpoint,
            "registered_tools": registered,
            "start": start_result,
        }

    _tool.__name__ = "builtin_mcp_register_server"
    return _tool


def _make_shell_tool(options: DynamicRunOptions, mode: str):
    def _tool(
        command: str,
        cwd: str | None = None,
        timeout_seconds: int | None = None,
    ) -> dict[str, Any]:
        return _run_shell_command(
            command=command,
            cwd=cwd,
            timeout_seconds=timeout_seconds,
            mode=mode,
            options=options,
        )

    _tool.__name__ = f"builtin_{mode}"
    return _tool


def _split_names(value: str) -> list[str]:
    return [
        item.strip()
        for chunk in value.split(",")
        for item in chunk.split()
        if item.strip()
    ]


def _run_shell_command(
    *,
    command: str,
    cwd: str | None,
    timeout_seconds: int | None,
    mode: str,
    options: DynamicRunOptions,
) -> dict[str, Any]:
    project_root = Path(options.project_root or Path.cwd()).expanduser().resolve()
    run_cwd = _resolve_cwd(project_root, cwd)
    timeout = _coerce_timeout(timeout_seconds, options.shell_timeout_seconds)

    if not run_cwd.is_dir():
        return {
            "ok": False,
            "returncode": 127,
            "exit_code": 127,
            "stdout": "",
            "stderr": (
                f"cwd does not exist: {run_cwd}. Create the directory first "
                "with files_write_text/shell_workspace_write, or use the "
                "project name from this flow's input."
            ),
            "command": command,
            "cwd": str(run_cwd),
            "mode": mode,
        }

    error = _validate_command(command, mode, options)
    if error:
        return {
            "ok": False,
            "returncode": 126,
            "exit_code": 126,
            "stdout": "",
            "stderr": error,
            "command": command,
            "cwd": str(run_cwd),
            "mode": mode,
        }

    completed = subprocess.run(
        ["/bin/bash", "-lc", command],
        cwd=run_cwd,
        capture_output=True,
        text=True,
        timeout=timeout,
        check=False,
    )
    max_output = max(500, options.shell_output_max_chars or _DEFAULT_MAX_OUTPUT_CHARS)
    return {
        "ok": completed.returncode == 0,
        "returncode": completed.returncode,
        "exit_code": completed.returncode,
        "stdout": _trim_output(completed.stdout, max_output),
        "stderr": _trim_output(completed.stderr, max_output),
        "truncated": (
            len(completed.stdout) > max_output or len(completed.stderr) > max_output
        ),
        "command": command,
        "cwd": str(run_cwd),
        "mode": mode,
    }


def _trim_output(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return (
        f"[truncated to last {max_chars} chars]\n"
        + text[-max_chars:]
    )


def _coerce_timeout(value: Any, default: int) -> float:
    if value is None:
        return float(default)
    try:
        return float(value)
    except (TypeError, ValueError):
        return float(default)


def _resolve_cwd(project_root: Path, cwd: str | None) -> Path:
    raw = Path(cwd).expanduser() if cwd else project_root
    resolved = raw.resolve() if raw.is_absolute() else (project_root / raw).resolve()
    try:
        resolved.relative_to(project_root)
    except ValueError as exc:
        raise ValueError(f"cwd must stay inside project_root: {resolved}") from exc
    return resolved


def _validate_command(
    command: str,
    mode: str,
    options: DynamicRunOptions,
) -> str | None:
    stripped = command.strip()
    if not stripped:
        return "command must not be empty"
    if any(token in stripped for token in _CONTROL_TOKENS):
        return "shell control operators are not allowed in builtin shell tools"

    try:
        parts = shlex.split(stripped)
    except ValueError as exc:
        return f"invalid shell command: {exc}"
    if not parts:
        return "command must not be empty"

    executable = parts[0]
    if mode == "readonly":
        return _validate_readonly(parts)
    if mode == "project_exec":
        if executable in {"python", "python3", "pytest", "uv", "npm", "pnpm", "yarn"}:
            return None
        return f"command {executable!r} is not allowed for project_exec"
    if mode == "workspace_write":
        if executable in {"mkdir", "touch", "cp", "mv"}:
            return None
        return f"command {executable!r} is not allowed for workspace_write"
    if mode == "install_dependency":
        return _validate_install(parts, options)
    return f"unknown shell mode: {mode}"


def _validate_readonly(parts: list[str]) -> str | None:
    executable = parts[0]
    if executable in {"pwd", "ls", "find", "rg", "grep", "cat", "sed", "awk", "head", "tail", "wc"}:
        return None
    if executable == "git":
        subcommand = parts[1] if len(parts) > 1 else ""
        if subcommand in {"status", "diff", "log", "show", "branch"}:
            return None
        return f"git subcommand {subcommand!r} is not read-only allowlisted"
    return f"command {executable!r} is not allowed for readonly"


def _validate_install(parts: list[str], options: DynamicRunOptions) -> str | None:
    policy = options.dependency_policy
    if not policy.allow_install:
        return "dependency installation is disabled by DynamicRunOptions"

    manager = parts[0]
    effective_manager = manager
    if manager in {"python", "python3"} and parts[1:4] == ["-m", "pip", "install"]:
        effective_manager = "pip"

    if effective_manager not in policy.allowed_managers:
        return f"dependency manager {manager!r} is not allowed"

    if effective_manager in {"pip", "uv"} and "install" not in parts and "add" not in parts:
        return f"{effective_manager} command must use install/add for dependency installation"
    if effective_manager in {"npm", "pnpm", "yarn"} and not any(
        token in parts for token in {"install", "add"}
    ):
        return f"{effective_manager} command must use install/add for dependency installation"

    requested = _extract_requested_packages(parts)
    if policy.allowed_packages:
        disallowed = [pkg for pkg in requested if pkg not in policy.allowed_packages]
        if disallowed:
            return f"packages are not in allowed_packages: {disallowed}"
    denied = [pkg for pkg in requested if pkg in policy.denied_packages]
    if denied:
        return f"packages are denied by dependency policy: {denied}"
    return None


def _extract_requested_packages(parts: list[str]) -> list[str]:
    packages: list[str] = []
    skip_next = False
    for token in parts[1:]:
        if skip_next:
            skip_next = False
            continue
        if token in {"install", "add", "-m", "pip"}:
            continue
        if token.startswith("-"):
            if token in {"-r", "--requirement"}:
                skip_next = True
            continue
        packages.append(token.split("==", 1)[0].split(">=", 1)[0])
    return packages
