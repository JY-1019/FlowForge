"""Tests for the dynamic flow generation feature.

Covers:
- @global_config(dynamic_flow=True) flag
- CompiledAgent.add_flow() — dynamic flow injection
- CompiledAgent.recompile() — full DAG rebuild
- DynamicFlowGenerator.compile_flow_code() — code → FlowMeta
- Meta-flow DAG injection when dynamic_flow=True
- add_flow_to_dag() — incremental compilation
"""
import pytest
from unittest.mock import AsyncMock, patch

from flowforge import (
    global_config, flow, task, step,
    FlowForge, CompileError,
)
from flowforge.annotations.decorators import _GLOBAL_ATTR, _FLOW_ATTR
from flowforge.doc.models import FlowDoc
from flowforge.dynamic import meta_flow as dynamic_meta_flow
from flowforge.planner import ExecutionPlan
from flowforge.schema.compiler import add_flow_to_dag
from flowforge.schema.dag import NodeType
from flowforge.types import AgentSkill, ClaudeSkill, FunctionTool, LLMConfig
from flowforge import DynamicRunOptions


# ---------------------------------------------------------------------------
# Fixtures: simple agent definitions at module level
# ---------------------------------------------------------------------------

@flow(name="alpha", prompt="first flow")
class AlphaFlow:
    @task(name="a_task", prompt="alpha task")
    class ATask:
        @step(order=1, prompt="alpha step")
        async def a_step(ctx):
            return {"result": "alpha"}


@flow(name="beta", prompt="second flow")
class BetaFlow:
    @task(name="b_task", prompt="beta task")
    class BTask:
        @step(order=1, prompt="beta step")
        async def b_step(ctx):
            return {"result": "beta"}


# Extra flow used to test add_flow()
@flow(name="gamma", prompt="third flow added dynamically")
class GammaFlow:
    @task(name="g_task", prompt="gamma task")
    class GTask:
        @step(order=1, prompt="gamma step")
        async def g_step(ctx):
            return {"result": "gamma"}


@global_config(prompt="test agent")
class BasicAgent:
    AlphaFlow = AlphaFlow
    BetaFlow = BetaFlow


@global_config(prompt="dynamic agent", dynamic_flow=True)
class DynamicAgent:
    AlphaFlow = AlphaFlow


@global_config(prompt="empty dynamic agent", dynamic_flow=True)
class EmptyDynamicAgent:
    pass


# Custom _dynamic_generator flow for precedence test
@flow(name="_dynamic_generator", prompt="my custom generator")
class CustomGeneratorFlow:
    @task(name="custom_gen", prompt="custom generation")
    class CustomGenTask:
        @step(order=1, prompt="custom step")
        async def custom_step(ctx):
            return {"custom": True}


@global_config(prompt="agent with custom gen", dynamic_flow=True)
class CustomAgent:
    AlphaFlow = AlphaFlow
    CustomGeneratorFlow = CustomGeneratorFlow


@flow(name="paper_report_pipeline", prompt="consume already-fetched paper payload and build report")
class PaperReportPipelineFlow:
    @task(name="pipeline", prompt="run report pipeline")
    class PipelineTask:
        @step(order=1, prompt="record pipeline execution")
        async def run_pipeline(ctx):
            seq = list(ctx.input.get("sequence", [])) if isinstance(ctx.input, dict) else []
            seq.append("pipeline")
            return {"sequence": seq}


@flow(name="fetch_trending_papers", prompt="fetch trending papers first")
class FetchTrendingPapersFlow:
    @task(name="fetch", prompt="fetch trending paper payload")
    class FetchTask:
        @step(order=1, prompt="fetch")
        async def do_fetch(ctx):
            return {"sequence": ["fetch"]}


@flow(name="_dynamic_generator", prompt="inject fetch flow")
class PartialGapGeneratorFlow:
    @task(name="inject", prompt="inject missing flow")
    class InjectTask:
        @step(order=1, prompt="inject missing flow")
        async def inject_flow(ctx):
            agent = ctx.shared_data["_compiled_agent"]
            node_id = agent.add_flow(FetchTrendingPapersFlow)
            return {
                "success": True,
                "dynamic_flow": "fetch_trending_papers",
                "flow_meta_name": "fetch_trending_papers",
                "node_id": node_id,
                "generated_code": "from flowforge import flow, task, step",
            }


@global_config(prompt="partial gap agent", dynamic_flow=True)
class PartialGapAgent:
    PaperReportPipelineFlow = PaperReportPipelineFlow
    PartialGapGeneratorFlow = PartialGapGeneratorFlow


def _fresh_basic():
    """Create a fresh BasicAgent to avoid test-to-test state leakage.

    add_flow() mutates global_meta.flows, so each test that calls it
    must use its own agent instance.
    """
    @global_config(prompt="test agent")
    class _Agent:
        AlphaFlow = AlphaFlow
        BetaFlow = BetaFlow
    return _Agent


# ---------------------------------------------------------------------------
# Test: dynamic_flow flag on GlobalMeta
# ---------------------------------------------------------------------------

class TestDynamicFlowFlag:

    def test_default_is_false(self):
        meta = getattr(BasicAgent, _GLOBAL_ATTR)
        assert meta.dynamic_flow is False

    def test_dynamic_flow_true(self):
        meta = getattr(DynamicAgent, _GLOBAL_ATTR)
        assert meta.dynamic_flow is True

    def test_internal_dynamic_generator_steps_have_generous_timeout(self):
        engine = FlowForge.compile(DynamicAgent)

        for node_id in [
            "global._dynamic_generator._generate_and_run.analyse_gap[1]",
            "global._dynamic_generator._generate_and_run.plan_workflow[2]",
            "global._dynamic_generator._generate_and_run.select_capability[3]",
            "global._dynamic_generator._generate_and_run.mcp_provision[4]",
            "global._dynamic_generator._generate_and_run.synthesise_and_inject[5]",
        ]:
            node = engine.dag.get_node(node_id)
            assert node is not None
            assert node.meta.timeout_seconds >= 300


# ---------------------------------------------------------------------------
# Test: CompiledAgent.add_flow()
# ---------------------------------------------------------------------------

class TestAddFlow:

    def test_add_flow_creates_new_node(self):
        engine = FlowForge.compile(_fresh_basic())
        initial_count = len(engine.dag.get_all_nodes())

        node_id = engine.add_flow(GammaFlow)

        assert node_id == "global.gamma"
        assert engine.dag.get_node("global.gamma") is not None
        assert len(engine.dag.get_all_nodes()) > initial_count

    def test_add_flow_subtree_is_complete(self):
        engine = FlowForge.compile(_fresh_basic())
        engine.add_flow(GammaFlow)

        # Should have: global.gamma, global.gamma.g_task,
        # global.gamma.g_task.g_step[1]
        assert engine.dag.get_node("global.gamma") is not None
        assert engine.dag.get_node("global.gamma.g_task") is not None
        assert engine.dag.get_node("global.gamma.g_task.g_step[1]") is not None

    def test_add_flow_duplicate_raises(self):
        engine = FlowForge.compile(_fresh_basic())
        with pytest.raises(CompileError, match="already exists"):
            engine.add_flow(AlphaFlow)

    def test_add_flow_not_decorated_raises(self):
        class Plain:
            pass

        engine = FlowForge.compile(_fresh_basic())
        with pytest.raises(CompileError, match="not decorated"):
            engine.add_flow(Plain)

    def test_add_flow_preserves_existing_nodes(self):
        engine = FlowForge.compile(_fresh_basic())
        # Capture existing nodes.
        original_ids = {n.id for n in engine.dag.get_all_nodes()}

        engine.add_flow(GammaFlow)

        current_ids = {n.id for n in engine.dag.get_all_nodes()}
        assert original_ids.issubset(current_ids)


# ---------------------------------------------------------------------------
# Test: CompiledAgent.recompile()
# ---------------------------------------------------------------------------

class TestRecompile:

    def test_recompile_rebuilds_dag(self):
        engine = FlowForge.compile(_fresh_basic())
        engine.add_flow(GammaFlow)
        count_after_add = len(engine.dag.get_all_nodes())

        engine.recompile()
        count_after_recompile = len(engine.dag.get_all_nodes())

        # After recompile, gamma is still present (it was added to global_meta.flows)
        assert count_after_recompile == count_after_add


# ---------------------------------------------------------------------------
# Test: add_flow_to_dag() — incremental compilation
# ---------------------------------------------------------------------------

class TestAddFlowToDag:

    def test_add_flow_to_existing_dag(self):
        engine = FlowForge.compile(_fresh_basic())
        dag = engine.dag

        meta = getattr(GammaFlow, _FLOW_ATTR)
        node_id = add_flow_to_dag(dag, meta)

        assert node_id == "global.gamma"
        node = dag.get_node(node_id)
        assert node is not None
        assert node.type == NodeType.FLOW

    def test_add_flow_duplicate_raises(self):
        engine = FlowForge.compile(_fresh_basic())
        dag = engine.dag

        meta = getattr(AlphaFlow, _FLOW_ATTR)
        with pytest.raises(CompileError, match="already exists"):
            add_flow_to_dag(dag, meta)


# ---------------------------------------------------------------------------
# Test: Meta-flow injection with dynamic_flow=True
# ---------------------------------------------------------------------------

class TestMetaFlowInjection:

    def test_dynamic_agent_has_meta_flow(self):
        engine = FlowForge.compile(DynamicAgent)

        # The _dynamic_generator flow should be present.
        node = engine.dag.get_node("global._dynamic_generator")
        assert node is not None
        assert node.type == NodeType.FLOW

    def test_dynamic_agent_has_meta_flow_steps(self):
        engine = FlowForge.compile(DynamicAgent)

        # Check the task exists.
        task_node = engine.dag.get_node(
            "global._dynamic_generator._generate_and_run"
        )
        assert task_node is not None

        # Check steps exist (5 steps: analyse_gap, plan_workflow,
        # select_capability, mcp_provision, synthesise_and_inject).
        step_ids = [
            n.id for n in engine.dag.get_all_nodes()
            if n.id.startswith("global._dynamic_generator._generate_and_run.")
            and n.type == NodeType.STEP
        ]
        assert len(step_ids) == 5

    def test_basic_agent_has_no_meta_flow(self):
        engine = FlowForge.compile(BasicAgent)
        assert engine.dag.get_node("global._dynamic_generator") is None

    def test_compile_idempotent(self):
        """Calling compile() twice does not duplicate the meta-flow."""
        engine1 = FlowForge.compile(DynamicAgent)
        engine2 = FlowForge.compile(DynamicAgent)

        gen_nodes_1 = [
            n for n in engine1.dag.get_all_nodes()
            if "_dynamic_generator" in n.id
        ]
        gen_nodes_2 = [
            n for n in engine2.dag.get_all_nodes()
            if "_dynamic_generator" in n.id
        ]
        assert len(gen_nodes_1) == len(gen_nodes_2)

    def test_user_custom_dynamic_generator_takes_precedence(self):
        """If user defines their own _dynamic_generator flow, use it."""
        engine = FlowForge.compile(CustomAgent)

        # The _dynamic_generator node should exist
        node = engine.dag.get_node("global._dynamic_generator")
        assert node is not None

        # It should be the USER's version (has "custom_gen" task,
        # not "_generate_and_run" from the built-in)
        task_node = engine.dag.get_node(
            "global._dynamic_generator.custom_gen"
        )
        assert task_node is not None

        # The built-in "_generate_and_run" task should NOT exist
        builtin_task = engine.dag.get_node(
            "global._dynamic_generator._generate_and_run"
        )
        assert builtin_task is None


# ---------------------------------------------------------------------------
# Test: DynamicFlowGenerator.compile_flow_code()
# ---------------------------------------------------------------------------

class TestCompileFlowCode:

    def test_compile_valid_code(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.types import LLMConfig

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(),
            dag=engine.dag,
        )

        code = '''
from flowforge import flow, task, step

@flow(name="hello_flow", prompt="says hello")
class HelloFlow:
    @task(name="greet", prompt="greet the user")
    class GreetTask:
        @step(order=1, prompt="say hello")
        async def say_hello(ctx):
            return {"message": "hello!"}
'''
        meta = generator.compile_flow_code(code)
        assert meta.name == "hello_flow"
        assert len(meta.tasks) == 1
        assert meta.tasks[0].name == "greet"

    def test_compile_invalid_code_raises(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.types import LLMConfig

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(),
            dag=engine.dag,
        )

        with pytest.raises(CompileError):
            generator.compile_flow_code("this is not python code!!!")

    def test_compile_code_without_flow_raises(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.types import LLMConfig

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(),
            dag=engine.dag,
        )

        code = '''
class JustAClass:
    pass
'''
        with pytest.raises(CompileError, match="does not contain"):
            generator.compile_flow_code(code)


# ---------------------------------------------------------------------------
# Test: helper functions
# ---------------------------------------------------------------------------

class TestHelpers:

    def test_strip_markdown_fences(self):
        from flowforge.dynamic.generator import _strip_markdown_fences

        assert _strip_markdown_fences("```python\ncode\n```") == "code"
        assert _strip_markdown_fences("```\ncode\n```") == "code"
        assert _strip_markdown_fences("code") == "code"

    def test_sanitise_name(self):
        from flowforge.dynamic.generator import _sanitise_name

        assert _sanitise_name("hello-world") == "hello_world"
        assert _sanitise_name("123abc") == "flow_123abc"
        assert _sanitise_name("My Flow!") == "my_flow"
        assert _sanitise_name("") == "dynamic_flow"
        assert _sanitise_name("valid_name") == "valid_name"

    @pytest.mark.asyncio
    async def test_gap_analysis_uses_heuristic_when_llm_fails(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        with patch(
            "flowforge.llm.caller.call_with_tool",
            new_callable=AsyncMock,
        ) as mock_call:
            mock_call.side_effect = TimeoutError("network timeout")
            result = await generator.analyse_gap(
                "세계에서 가장 높은 산 Top 5를 표로 정리해줘"
            )

        assert result["covered"] is False
        assert result["suggested_flow_name"] == "top5_highest_mountains_table"
        assert "heuristic" in result["reason"].lower()


class TestDynamicGeneratorPrompting:

    @pytest.mark.asyncio
    async def test_tool_catalog_included_in_codegen_prompt(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(BasicAgent)
        tool = FunctionTool(
            func=lambda limit=3: {"limit": limit},
            name="fetch_hf_trending_papers",
            description="Fetch top trending papers from Hugging Face.",
        )
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
            tool_configs=[tool],
        )

        with patch("flowforge.execution.llm.call_llm_api", new_callable=AsyncMock) as mock_api:
            mock_api.return_value = "from flowforge import flow, task, step"
            await generator.generate_flow_code(
                flow_name="fetch_trending_papers",
                flow_prompt="fetch trending papers",
                user_query="latest papers",
                project_context="pyproject uses pytest and flowforge decorators",
            )

            user_prompt = mock_api.call_args.kwargs["user_prompt"]
            call_kwargs = mock_api.call_args.kwargs
            assert "Available tools:" in user_prompt
            assert "fetch_hf_trending_papers" in user_prompt
            assert "Fetch top trending papers" in user_prompt
            assert 'Decorator scope: tools=["fetch_hf_trending_papers"]' in user_prompt
            assert "Code generation context prepared" in user_prompt
            assert "pyproject uses pytest" in user_prompt
            assert call_kwargs["tool_configs"] is None

    @pytest.mark.asyncio
    async def test_codegen_tool_use_can_be_enabled_explicitly(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(BasicAgent)
        tool = FunctionTool(
            func=lambda limit=3: {"limit": limit},
            name="fetch_hf_trending_papers",
            description="Fetch top trending papers from Hugging Face.",
        )
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
            tool_configs=[tool],
            dynamic_options=DynamicRunOptions(allow_codegen_tool_use=True),
        )

        with patch("flowforge.execution.llm.call_llm_api", new_callable=AsyncMock) as mock_api:
            mock_api.return_value = "from flowforge import flow, task, step"
            await generator.generate_flow_code(
                flow_name="fetch_trending_papers",
                flow_prompt="fetch trending papers",
                user_query="latest papers",
            )

            assert mock_api.call_args.kwargs["tool_configs"] == [tool]

    @pytest.mark.asyncio
    async def test_codegen_includes_prompt_only_skills_by_default(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        skill_dir = tmp_path / "skills" / "clone-coding"
        skill_dir.mkdir(parents=True)
        (skill_dir / "SKILL.md").write_text(
            "---\n"
            "name: clone-coding\n"
            "description: Clone-coding guidance.\n"
            "---\n"
            "# Clone Coding\n",
            encoding="utf-8",
        )
        executable_tool = FunctionTool(
            func=lambda: {"ok": True},
            name="executable_tool",
        )
        claude_skill = ClaudeSkill(name="frontend-design")
        agent_skill = AgentSkill(path=str(skill_dir))
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig.for_claude(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[executable_tool, claude_skill, agent_skill],
            dynamic_options=DynamicRunOptions(allow_codegen_tool_use=False),
        )

        with patch("flowforge.execution.llm.call_llm_api", new_callable=AsyncMock) as mock_api:
            mock_api.return_value = "from flowforge import flow, task, step"
            await generator.generate_flow_code(
                flow_name="clone_flow",
                flow_prompt="clone a frontend",
                user_query="clone naver",
            )

            assert mock_api.call_args.kwargs["tool_configs"] == [
                claude_skill,
                agent_skill,
            ]

    @pytest.mark.asyncio
    async def test_codegen_system_requires_prompts_and_tool_refs(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        tool = FunctionTool(
            func=lambda url: {"ok": True, "url": url},
            name="web_fetch_url",
            description="Fetch a public web page.",
        )
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[tool],
        )

        with patch("flowforge.execution.llm.call_llm_api", new_callable=AsyncMock) as mock_api:
            mock_api.return_value = "from flowforge import flow, task, step"
            await generator.generate_flow_code(
                flow_name="clone_site",
                flow_prompt="clone a public site",
                user_query={
                    "target_url": "https://example.com",
                    "required_tools": ["web_fetch_url"],
                },
            )

            system_prompt = mock_api.call_args.kwargs["system_prompt"]
            context_prompt = generator.build_code_generation_context(
                flow_name="clone_site",
                flow_prompt="clone a public site",
                user_query="clone https://example.com",
            )

        assert "Every generated `@flow`, `@task`, and `@step` MUST have" in system_prompt
        assert 'tools=["web_fetch_url"]' in system_prompt
        assert "<tool_name>" in system_prompt
        assert "Give every @task and @step a specific prompt" in context_prompt
        assert 'tools=["tool_name"]' in context_prompt

    @pytest.mark.asyncio
    async def test_generate_and_compile_retries_missing_required_tool_usage(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        tool = FunctionTool(
            func=lambda url: {"ok": True, "url": url},
            name="web_fetch_url",
            description="Fetch a public web page.",
        )
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[tool],
        )

        bad_code = '''
from flowforge import flow, task, step

@flow(name="fetch_page", prompt="Fetch a public page")
class FetchPageFlow:
    @task(name="fetch", prompt="Fetch the page")
    class FetchTask:
        @step(order=1, prompt="Return a placeholder")
        async def fetch(ctx):
            return {"ok": True}
'''
        fixed_code = '''
from flowforge import flow, task, step

@flow(name="fetch_page", prompt="Fetch a public page", tools=["web_fetch_url"])
class FetchPageFlow:
    @task(name="fetch", prompt="Fetch the page with the web tool", tools=["web_fetch_url"])
    class FetchTask:
        @step(order=1, prompt="Call web_fetch_url for the target URL", tools=["web_fetch_url"])
        async def fetch(ctx):
            target_url = ctx.input.get("target_url", "https://example.com")
            return await ctx.call_tool("web_fetch_url", url=target_url)
'''
        with patch("flowforge.execution.llm.call_llm_api", new_callable=AsyncMock) as mock_api:
            mock_api.side_effect = [bad_code, fixed_code]
            meta, code = await generator.generate_and_compile(
                flow_name="fetch_page",
                flow_prompt="Fetch a public page",
                user_query={
                    "target_url": "https://example.com",
                    "required_tools": ["web_fetch_url"],
                },
            )

        assert meta.name == "fetch_page"
        assert code == fixed_code.strip()
        assert mock_api.await_count == 2

    def test_check_tool_ref_validity_accepts_registered_names(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        tool = FunctionTool(
            func=lambda url: {"ok": True, "url": url},
            name="web_fetch_url",
            description="Fetch a public web page.",
        )
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[tool],
        )

        good_code = '''
from flowforge import flow, task, step

@flow(name="fetch_page", prompt="Fetch a public page", tools=["web_fetch_url"])
class FetchPageFlow:
    @task(name="fetch", prompt="Fetch the page", tools=["web_fetch_url"])
    class FetchTask:
        @step(order=1, prompt="Use <web_fetch_url>", tools=["web_fetch_url"])
        async def fetch(ctx):
            return await ctx.call_tool("web_fetch_url", url="https://example.com")
'''
        assert generator.check_tool_ref_validity(good_code) is None

    def test_check_tool_ref_validity_rejects_unknown_names(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        skill_dir = tmp_path / "skills" / "frontend-design"
        skill_dir.mkdir(parents=True)
        (skill_dir / "SKILL.md").write_text(
            "---\nname: frontend-design\ndescription: design.\n---\n# guidance\n",
            encoding="utf-8",
        )
        agent_skill = AgentSkill(path=str(skill_dir), name="frontend-design")
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[agent_skill],
        )

        hallucinated_code = '''
from flowforge import flow, task, step

@flow(name="clone_site", prompt="Clone a site", tools=["frontend-design"])
class CloneFlow:
    @task(name="design", prompt="Design step", tools=["frontend-design"])
    class DesignTask:
        @step(order=1, prompt="Apply <clone-coding> to draft markup")
        async def draft(ctx):
            return await ctx.call_llm("Use <clone-coding> to render")
'''
        error = generator.check_tool_ref_validity(hallucinated_code)
        assert error is not None
        assert "clone-coding" in error
        assert "frontend-design" in error

    def test_check_tool_ref_validity_skipped_when_no_tools(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(BasicAgent).dag,
            tool_configs=[],
        )

        any_code = '''
from flowforge import flow, task, step

@flow(name="x", prompt="x")
class X:
    @task(name="t", prompt="t")
    class T:
        @step(order=1, prompt="call <unknown>")
        async def s(ctx):
            return {}
'''
        assert generator.check_tool_ref_validity(any_code) is None

    @pytest.mark.asyncio
    async def test_precomputed_gap_skips_gap_analysis(self):
        engine = FlowForge.compile(DynamicAgent)

        class _Ctx:
            input = {
                "user_query": "need missing flow",
                "gap_analysis": {
                    "covered": False,
                    "reason": "missing fetch flow",
                    "suggested_flow_name": "fetch_trending_papers",
                    "suggested_flow_prompt": "fetch trending papers first",
                },
            }

            class _Global:
                llm_config = LLMConfig()

            global_ctx = _Global()
            shared_data = {"_compiled_agent": engine}

        with patch(
            "flowforge.dynamic.generator.DynamicFlowGenerator.analyse_gap",
            new_callable=AsyncMock,
        ) as mock_analyse:
            result = await dynamic_meta_flow._analyse_gap(_Ctx())

        mock_analyse.assert_not_called()
        assert result["precomputed_gap"] is True
        assert result["flow_name"] == "fetch_trending_papers"

class TestDynamicRunOptionsAndManifest:
    def test_compile_with_dynamic_options_adds_builtin_shell_tools(self, tmp_path):
        options = DynamicRunOptions(project_root=str(tmp_path))

        @global_config(prompt="dynamic builtin agent", dynamic_flow=True)
        class _Agent:
            AlphaFlow = AlphaFlow

        engine = FlowForge.compile(_Agent, dynamic_options=options)

        names = [
            tool.name or getattr(tool.func, "__name__", "")
            for tool in engine._global_meta.tools
            if isinstance(tool, FunctionTool)
        ]
        assert "shell_readonly" in names
        assert "shell_project_exec" in names
        assert "shell_workspace_write" not in names

    def test_compile_autoloads_generated_flow_manifest(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            generated_step_timeout_seconds=240,
        )
        persist_flow_code(
            flow_name="persisted_flow",
            code="""
from flowforge import flow, task, step

@flow(name="persisted_flow", prompt="loaded from manifest")
class PersistedFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="run")
        async def run(ctx):
            return {"loaded": True}
""",
            options=options,
        )

        @global_config(prompt="manifest agent", dynamic_flow=True)
        class _Agent:
            AlphaFlow = AlphaFlow

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        node = engine.dag.get_node("global.persisted_flow.main.run[1]")
        assert node is not None
        assert node.meta.timeout_seconds == 240
        flow_node = engine.dag.get_node("global.persisted_flow")
        assert flow_node is not None
        assert getattr(flow_node.meta, "__flowforge_dynamic_source_code__", "")

    @pytest.mark.asyncio
    async def test_manifest_bridge_flow_runs_before_downstream(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code

        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(
            flow_name="manifest_fetch",
            code="""
from flowforge import flow, task, step

@flow(name="manifest_fetch", prompt="fetch before report")
class ManifestFetchFlow:
    @task(name="fetch", prompt="fetch")
    class FetchTask:
        @step(order=1, prompt="fetch")
        async def fetch(ctx):
            return {"sequence": ["manifest_fetch"]}
""",
            options=options,
            inject_before="paper_report_pipeline",
            downstream_flow_route="paper_report_pipeline",
        )

        @global_config(prompt="manifest bridge agent", dynamic_flow=True)
        class _Agent:
            PaperReportPipelineFlow = PaperReportPipelineFlow

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        roots = [
            node.name for node in engine.dag.get_children("global")
            if node.type == NodeType.FLOW
        ]
        assert roots.index("manifest_fetch") < roots.index("paper_report_pipeline")

        result = await engine.run(
            {"sequence": []},
            route=["paper_report_pipeline", "manifest_fetch"],
        )
        assert result["sequence"] == ["manifest_fetch", "pipeline"]

    def test_compile_autoloads_generated_tool_manifest(self, tmp_path):
        from flowforge.dynamic.manifest import persist_tool_code

        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_tool_code(
            tool_name="project_echo",
            code="""
def project_echo(text: str) -> dict:
    return {"text": text}
""",
            options=options,
        )

        @global_config(prompt="manifest tool agent", dynamic_flow=True)
        class _Agent:
            AlphaFlow = AlphaFlow

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        names = [
            tool.name or getattr(tool.func, "__name__", "")
            for tool in engine._global_meta.tools
            if isinstance(tool, FunctionTool)
        ]
        assert "project_echo" in names

    @pytest.mark.asyncio
    async def test_autoloaded_generated_flow_repairs_and_persists(
        self, tmp_path
    ):
        from flowforge.dynamic.manifest import persist_flow_code

        broken_code = """
from flowforge import flow, task, step

@flow(name="repairable_flow", prompt="loaded broken flow")
class RepairableFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="break")
        async def run(ctx):
            raise RuntimeError("boom from persisted flow")
"""
        fixed_code = """
from flowforge import flow, task, step

@flow(name="repairable_flow", prompt="fixed flow")
class RepairableFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="fixed")
        async def run(ctx):
            return {"fixed_marker": True}
"""
        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(
            flow_name="repairable_flow",
            code=broken_code,
            options=options,
            class_name="RepairableFlow",
        )

        @global_config(prompt="manifest repair agent", dynamic_flow=True)
        class _Agent:
            pass

        async def fake_fix(self, code, error, user_query, **kwargs):
            assert "boom from persisted flow" in error
            return fixed_code

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        with patch(
            "flowforge.dynamic.generator.DynamicFlowGenerator.fix_code",
            new=fake_fix,
        ):
            result = await engine.run({}, route="repairable_flow")

        assert result == {"fixed_marker": True}
        generated_file = (
            tmp_path
            / "flowforge"
            / "generated"
            / "flows"
            / "repairable_flow.py"
        )
        assert "fixed_marker" in generated_file.read_text(encoding="utf-8")

    @pytest.mark.asyncio
    async def test_autoloaded_generated_flow_preflight_repairs_quality(
        self, tmp_path
    ):
        from flowforge.dynamic.manifest import persist_flow_code

        stale_code = """
from flowforge import flow, task, step

@flow(name="stale_quality_flow", prompt="loaded stale flow")
class StaleQualityFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="return placeholder")
        async def run(ctx):
            return {"title": "HTML 본문 데이터 없음"}
"""
        fixed_code = """
from flowforge import flow, task, step

@flow(name="stale_quality_flow", prompt="fixed flow")
class StaleQualityFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="fixed")
        async def run(ctx):
            return {"title": "real story"}
"""
        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(
            flow_name="stale_quality_flow",
            code=stale_code,
            options=options,
            class_name="StaleQualityFlow",
        )

        @global_config(prompt="manifest quality repair agent", dynamic_flow=True)
        class _Agent:
            pass

        async def fake_fix(self, code, error, user_query, **kwargs):
            assert "Preflight generated-code quality check failed" in error
            return fixed_code

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        with patch(
            "flowforge.dynamic.generator.DynamicFlowGenerator.fix_code",
            new=fake_fix,
        ):
            result = await engine.run({}, route="stale_quality_flow")

        assert result == {"title": "real story"}

    @pytest.mark.asyncio
    async def test_generated_flow_repairs_placeholder_runtime_output(
        self, tmp_path
    ):
        from flowforge.dynamic.manifest import persist_flow_code

        runtime_bad_code = """
from flowforge import flow, task, step

@flow(name="runtime_placeholder_flow", prompt="runtime placeholder flow")
class RuntimePlaceholderFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="echo title")
        async def run(ctx):
            return {"title": ctx.input["title"]}
"""
        fixed_code = """
from flowforge import flow, task, step

@flow(name="runtime_placeholder_flow", prompt="fixed flow")
class RuntimePlaceholderFlow:
    @task(name="main", prompt="main")
    class MainTask:
        @step(order=1, prompt="fixed")
        async def run(ctx):
            return {"title": "real story"}
"""
        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(
            flow_name="runtime_placeholder_flow",
            code=runtime_bad_code,
            options=options,
            class_name="RuntimePlaceholderFlow",
        )

        @global_config(prompt="runtime quality repair agent", dynamic_flow=True)
        class _Agent:
            pass

        async def fake_fix(self, code, error, user_query, **kwargs):
            assert "placeholder failure output" in error
            return fixed_code

        engine = FlowForge.compile(_Agent, dynamic_options=options)
        with patch(
            "flowforge.dynamic.generator.DynamicFlowGenerator.fix_code",
            new=fake_fix,
        ):
            result = await engine.run(
                {"title": "HTML 잘림으로 확인 불가"},
                route="runtime_placeholder_flow",
            )

        assert result == {"title": "real story"}


class TestPartialGapDynamicExecution:

    @pytest.mark.asyncio
    async def test_partial_gap_triggers_single_dynamic_replan(self):
        engine = FlowForge.compile(PartialGapAgent)
        engine.docs["global.paper_report_pipeline"] = FlowDoc(
            summary="Build the report once the paper payload already exists.",
            preconditions=["Paper payload must be fetched first."],
        )

        async def _plan_side_effect(user_request, dag, docs, llm_config):
            if dag.get_node("global.fetch_trending_papers") is None:
                return ExecutionPlan(
                    node_ids=sorted(dag.resolve_route("paper_report_pipeline")),
                    mode="autonomous",
                    rationale="pipeline exists but fetch flow is missing",
                    metadata={
                        "routes": ["paper_report_pipeline"],
                        "gap_detected": True,
                        "reason": "Need a fetch flow before the report pipeline can run.",
                        "suggested_flow_name": "fetch_trending_papers",
                        "suggested_flow_prompt": "Fetch trending papers and return paper payload JSON.",
                    },
                )

            return ExecutionPlan(
                node_ids=sorted(
                    dag.resolve_route("fetch_trending_papers")
                    | dag.resolve_route("paper_report_pipeline")
                ),
                mode="autonomous",
                rationale="fetch first, then pipeline",
                metadata={
                    "routes": ["fetch_trending_papers", "paper_report_pipeline"],
                    "gap_detected": False,
                    "reason": "",
                    "suggested_flow_name": "",
                    "suggested_flow_prompt": "",
                },
            )

        mock_plan = AsyncMock(side_effect=_plan_side_effect)
        with patch(
            "flowforge.planner.llm_planner.LLMPlanner.plan",
            new=mock_plan,
        ):
            result, trace = await engine.run_traced(
                "Make the trending papers report",
                planning_mode="autonomous",
            )

        assert mock_plan.call_count == 2
        assert result["sequence"] == ["fetch", "pipeline"]
        assert engine.last_dynamic_generation["dynamic_flow"] == "fetch_trending_papers"
        assert trace.get_node_trace("global._dynamic_generator") is not None

        flow_visits = [
            node.name for node in trace.nodes_by_type("flow")
            if node.node_id in {
                "global._dynamic_generator",
                "global.fetch_trending_papers",
                "global.paper_report_pipeline",
            }
        ]
        assert flow_visits == [
            "_dynamic_generator",
            "fetch_trending_papers",
            "paper_report_pipeline",
        ]


# ---------------------------------------------------------------------------
# Test: end-to-end add_flow + run
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_add_flow_and_run():
    """After add_flow(), the new flow is runnable via route."""
    engine = FlowForge.compile(_fresh_basic())
    engine.add_flow(GammaFlow)

    result = await engine.run({"input": "test"}, route="gamma")
    assert result == {"result": "gamma"}


@pytest.mark.asyncio
async def test_original_flows_still_work_after_add():
    """Adding a flow does not break existing flows."""
    engine = FlowForge.compile(_fresh_basic())
    engine.add_flow(GammaFlow)

    result = await engine.run({"input": "test"}, route="alpha")
    assert result == {"result": "alpha"}


# ---------------------------------------------------------------------------
# Test: schema-aware dynamic flow generation (downstream contract)
# ---------------------------------------------------------------------------

from pydantic import BaseModel   # noqa: E402  (used by fixtures below)


class _PaperPayload(BaseModel):
    papers: list[dict]
    fetched_at: str


@flow(
    name="paper_pipeline_with_schema",
    prompt="consume a strongly-typed paper payload and build report",
)
class PaperPipelineWithSchemaFlow:
    @task(
        name="pipeline",
        prompt="process paper payload",
        input_schema=_PaperPayload,
    )
    class PipelineTask:
        @step(
            order=1,
            prompt="process",
            input_schema=_PaperPayload,
        )
        async def run_pipeline(ctx):
            payload = ctx.input
            return {"papers": payload.papers, "processed": True}


@global_config(prompt="schema-aware agent", dynamic_flow=True)
class SchemaAwareAgent:
    PaperPipelineWithSchemaFlow = PaperPipelineWithSchemaFlow


class TestSchemaContract:
    """Dynamic codegen receives the downstream flow's input_schema as contract."""

    def test_entry_input_schema_helper_finds_first_step_schema(self):
        from flowforge.dynamic.generator import _entry_input_schema
        from flowforge.annotations.decorators import _FLOW_ATTR

        flow_meta = getattr(PaperPipelineWithSchemaFlow, _FLOW_ATTR)
        assert _entry_input_schema(flow_meta) is _PaperPayload

    def test_entry_input_schema_returns_none_when_nothing_declared(self):
        from flowforge.dynamic.generator import _entry_input_schema
        from flowforge.annotations.decorators import _FLOW_ATTR

        flow_meta = getattr(AlphaFlow, _FLOW_ATTR)
        assert _entry_input_schema(flow_meta) is None

    def test_resolve_downstream_contract_returns_json_schema(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(SchemaAwareAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        contract, downstream_name = generator.resolve_downstream_contract(
            "paper_pipeline_with_schema",
        )
        assert downstream_name == "paper_pipeline_with_schema"
        assert contract is not None
        assert set(contract.get("required", [])) == {"papers", "fetched_at"}
        assert "papers" in contract["properties"]

    def test_resolve_downstream_contract_returns_none_for_unknown_route(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(SchemaAwareAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        contract, name = generator.resolve_downstream_contract("does_not_exist")
        assert contract is None
        assert name is None

    @pytest.mark.asyncio
    async def test_codegen_prompt_includes_contract_json_schema(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(SchemaAwareAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )
        contract, downstream_name = generator.resolve_downstream_contract(
            "paper_pipeline_with_schema",
        )

        with patch(
            "flowforge.execution.llm.call_llm_api",
            new_callable=AsyncMock,
        ) as mock_api:
            mock_api.return_value = "from flowforge import flow, task, step"

            await generator.generate_flow_code(
                flow_name="fetch_trending_papers",
                flow_prompt="fetch trending papers",
                user_query="give me the trending papers",
                downstream_contract=contract,
                downstream_flow_name=downstream_name,
            )

            user_prompt = mock_api.call_args.kwargs["user_prompt"]
            assert "Downstream input contract" in user_prompt
            assert "paper_pipeline_with_schema" in user_prompt
            assert "\"fetched_at\"" in user_prompt
            assert "\"papers\"" in user_prompt

    def test_compatibility_check_flags_missing_required_keys(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(SchemaAwareAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        class _Wrong(BaseModel):
            unrelated: str

        @flow(name="fetch_bad", prompt="bad fetch")
        class _BadFetchFlow:
            @task(name="fetch", prompt="fetch", output_schema=_Wrong)
            class _Fetch:
                @step(order=1, prompt="fetch", output_schema=_Wrong)
                async def do_fetch(ctx):
                    return {"unrelated": "oops"}

        from flowforge.annotations.decorators import _FLOW_ATTR
        flow_meta = getattr(_BadFetchFlow, _FLOW_ATTR)

        contract = _PaperPayload.model_json_schema()
        err = generator.check_contract_compatibility(flow_meta, contract)
        assert err is not None
        assert "papers" in err or "fetched_at" in err

    def test_compatibility_check_passes_when_keys_match(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(SchemaAwareAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        @flow(name="fetch_ok", prompt="ok fetch")
        class _OkFetchFlow:
            @task(name="fetch", prompt="fetch", output_schema=_PaperPayload)
            class _Fetch:
                @step(order=1, prompt="fetch", output_schema=_PaperPayload)
                async def do_fetch(ctx):
                    return {"papers": [], "fetched_at": "2026-01-01"}

        from flowforge.annotations.decorators import _FLOW_ATTR
        flow_meta = getattr(_OkFetchFlow, _FLOW_ATTR)

        contract = _PaperPayload.model_json_schema()
        err = generator.check_contract_compatibility(flow_meta, contract)
        assert err is None

    def test_build_dynamic_input_forwards_downstream_route(self):
        engine = FlowForge.compile(SchemaAwareAgent)

        class _Plan:
            metadata = {
                "gap_detected": True,
                "suggested_flow_name": "fetch_trending_papers",
                "suggested_flow_prompt": "Fetch trending papers.",
                "reason": "missing fetch flow",
                "downstream_flow_route": "paper_pipeline_with_schema",
            }

        payload = engine._engine._build_dynamic_input(
            "give me trending papers", _Plan(),
        )
        assert isinstance(payload, dict)
        assert payload["downstream_flow_route"] == "paper_pipeline_with_schema"
        assert payload["gap_analysis"]["suggested_flow_name"] == "fetch_trending_papers"

    def test_build_dynamic_input_synthesizes_missing_gap_fields(self):
        engine = FlowForge.compile(EmptyDynamicAgent)

        class _Plan:
            metadata = {
                "gap_detected": True,
                "reason": "Need a clone coding flow for a frontend project.",
                "suggested_flow_name": "",
                "suggested_flow_prompt": "",
                "downstream_flow_route": "",
            }

        payload = engine._engine._build_dynamic_input(
            "Create a Naver clone project", _Plan(),
        )

        assert isinstance(payload, dict)
        assert payload["gap_analysis"]["suggested_flow_name"].startswith("dynamic_")
        assert "clone coding flow" in payload["gap_analysis"]["suggested_flow_prompt"]

    def test_build_runtime_failure_dynamic_input_targets_downstream(self):
        engine = FlowForge.compile(DynamicAgent)
        payload = engine._engine._build_runtime_failure_dynamic_input(
            input_data="fetch arXiv papers and make a PPT",
            error=ValueError("Could not parse JSON-like paper payload."),
            downstream_route="paper_report_pipeline",
        )

        assert payload["gap_analysis"]["covered"] is False
        assert payload["downstream_flow_route"] == "paper_report_pipeline"
        assert "paper_report_pipeline" in payload["gap_analysis"]["reason"]
        assert "upstream" in payload["gap_analysis"]["suggested_flow_prompt"]

    @pytest.mark.asyncio
    async def test_meta_flow_analyse_gap_propagates_downstream_route(self):
        engine = FlowForge.compile(SchemaAwareAgent)

        class _Ctx:
            input = {
                "user_query": "need missing flow",
                "gap_analysis": {
                    "covered": False,
                    "reason": "missing fetch flow",
                    "suggested_flow_name": "fetch_trending_papers",
                    "suggested_flow_prompt": "fetch",
                },
                "downstream_flow_route": "paper_pipeline_with_schema",
            }

            class _Global:
                llm_config = LLMConfig()

            global_ctx = _Global()
            shared_data = {"_compiled_agent": engine}

        result = await dynamic_meta_flow._analyse_gap(_Ctx())
        assert result["downstream_flow_route"] == "paper_pipeline_with_schema"
        assert result["precomputed_gap"] is True


# ---------------------------------------------------------------------------
# Test: AST safety validation
# ---------------------------------------------------------------------------

class TestASTSafety:

    def test_safe_code_passes(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = '''
from flowforge import flow, task, step

@flow(name="ok", prompt="ok")
class OkFlow:
    @task(name="t", prompt="t")
    class T:
        @step(order=1, prompt="s")
        async def s(ctx):
            return {"result": "ok"}
'''
        assert _validate_generated_ast(code) is None

    def test_os_system_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = 'import os\nos.system("rm -rf /")'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "os.system" in error

    def test_subprocess_import_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = 'import subprocess\nsubprocess.run(["ls"])'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "subprocess" in error

    def test_eval_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = 'eval("1+1")'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "eval" in error

    def test_exec_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = 'exec("print(1)")'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "exec" in error

    def test_dunder_import_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = '__import__("os")'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "__import__" in error

    def test_shutil_rmtree_rejected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        code = 'import shutil\nshutil.rmtree("/tmp")'
        error = _validate_generated_ast(code)
        assert error is not None
        assert "shutil" in error

    def test_syntax_error_detected(self):
        from flowforge.dynamic.generator import _validate_generated_ast

        error = _validate_generated_ast("def broken(")
        assert error is not None
        assert "SyntaxError" in error

    def test_compile_flow_code_rejects_unsafe_code(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(),
            dag=engine.dag,
        )

        unsafe_code = '''
import subprocess
from flowforge import flow, task, step

@flow(name="bad", prompt="bad")
class BadFlow:
    @task(name="t", prompt="t")
    class T:
        @step(order=1, prompt="s")
        async def s(ctx):
            subprocess.run(["ls"])
            return {"result": "bad"}
'''
        with pytest.raises(CompileError, match="safety check"):
            generator.compile_flow_code(unsafe_code)


# ---------------------------------------------------------------------------
# Test: generate_and_compile returns (meta, code) tuple
# ---------------------------------------------------------------------------

class TestGenerateAndCompileReturnType:

    @pytest.mark.asyncio
    async def test_generate_and_compile_returns_code(self):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        engine = FlowForge.compile(BasicAgent)
        generator = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=engine.dag,
        )

        generated_code = '''
from flowforge import flow, task, step

@flow(name="test_flow", prompt="test")
class TestFlowGenerated:
    @task(name="t", prompt="t")
    class T:
        @step(order=1, prompt="s")
        async def s(ctx):
            return {"result": "ok"}
'''
        with patch(
            "flowforge.execution.llm.call_llm_api",
            new_callable=AsyncMock,
        ) as mock_api:
            mock_api.return_value = generated_code
            meta, code = await generator.generate_and_compile(
                flow_name="test_flow",
                flow_prompt="test",
                user_query="test query",
            )

        assert meta.name == "test_flow"
        assert "TestFlowGenerated" in code


# ---------------------------------------------------------------------------
# Test: Builtin utility tools
# ---------------------------------------------------------------------------

class TestBuiltinUtilityTools:

    def test_utility_tools_included_in_pack(self, tmp_path):
        options = DynamicRunOptions(project_root=str(tmp_path))
        from flowforge.tools.builtin import create_builtin_tool_pack

        tools = create_builtin_tool_pack(options)
        names = [t.name for t in tools]

        assert "web_fetch_url" in names
        assert "json_select_fields" in names
        assert "files_read_text" in names
        assert "files_write_text" in names
        assert "files_list_dir" in names

    def test_json_select_fields(self, tmp_path):
        from flowforge.tools.builtin import _make_json_select_fields_tool

        tool = _make_json_select_fields_tool()
        result = tool(data='{"a": 1, "b": 2, "c": 3}', fields="a,c")
        assert result["ok"] is True
        assert result["selected"] == {"a": 1, "c": 3}
        assert result["missing"] == []

    def test_json_select_fields_missing_keys(self, tmp_path):
        from flowforge.tools.builtin import _make_json_select_fields_tool

        tool = _make_json_select_fields_tool()
        result = tool(data='{"a": 1}', fields="a,z")
        assert result["ok"] is True
        assert result["missing"] == ["z"]

    def test_files_read_write_roundtrip(self, tmp_path):
        from flowforge.tools.builtin import (
            _make_files_read_text_tool,
            _make_files_write_text_tool,
        )

        write_tool = _make_files_write_text_tool(tmp_path)
        read_tool = _make_files_read_text_tool(tmp_path, 4000)

        write_result = write_tool(path="test.txt", content="hello world")
        assert write_result["ok"] is True

        read_result = read_tool(path="test.txt")
        assert read_result["ok"] is True
        assert read_result["content"] == "hello world"

    def test_files_list_dir(self, tmp_path):
        from flowforge.tools.builtin import _make_files_list_dir_tool

        (tmp_path / "a.txt").write_text("a")
        (tmp_path / "b.txt").write_text("b")
        (tmp_path / "subdir").mkdir()

        tool = _make_files_list_dir_tool(tmp_path)
        result = tool(path=".")
        assert result["ok"] is True
        names = [e["name"] for e in result["entries"]]
        assert "a.txt" in names
        assert "subdir" in names

    def test_files_read_rejects_path_outside_project(self, tmp_path):
        from flowforge.tools.builtin import _make_files_read_text_tool

        tool = _make_files_read_text_tool(tmp_path, 4000)
        result = tool(path="/etc/passwd")
        assert result["ok"] is False
        assert "project root" in result["error"]

    def test_pip_install_tool_included_in_pack(self, tmp_path):
        options = DynamicRunOptions(project_root=str(tmp_path))
        from flowforge.tools.builtin import create_builtin_tool_pack

        tools = create_builtin_tool_pack(options)
        names = [t.name for t in tools]
        assert "pip_install" in names

    def test_pip_install_blocked_when_disabled(self, tmp_path):
        from flowforge.tools.builtin import _make_pip_install_tool
        from flowforge.types import DependencyPolicy

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            dependency_policy=DependencyPolicy(allow_install=False),
        )
        tool = _make_pip_install_tool(options)
        result = tool(packages="some-package")
        assert result["ok"] is False
        assert "disabled" in result["error"]

    def test_pip_install_denied_package(self, tmp_path):
        from flowforge.tools.builtin import _make_pip_install_tool
        from flowforge.types import DependencyPolicy

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            dependency_policy=DependencyPolicy(
                allow_install=True,
                denied_packages=["evil-package"],
            ),
        )
        tool = _make_pip_install_tool(options)
        result = tool(packages="evil-package")
        assert result["ok"] is False
        assert "denied" in result["error"].lower()

    def test_pip_install_allowed_packages_filter(self, tmp_path):
        from flowforge.tools.builtin import _make_pip_install_tool
        from flowforge.types import DependencyPolicy

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            dependency_policy=DependencyPolicy(
                allow_install=True,
                allowed_packages=["safe-pkg"],
            ),
        )
        tool = _make_pip_install_tool(options)
        result = tool(packages="not-in-allowlist")
        assert result["ok"] is False
        assert "allowed_packages" in result["error"]

    # ── Document tools ────────────────────────────────────────────────

    def test_document_tools_included_in_pack(self, tmp_path):
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(project_root=str(tmp_path))
        tools = create_builtin_tool_pack(options)
        names = [t.name for t in tools]

        assert "pdf_read_text" in names
        assert "pptx_create" in names
        assert "csv_read" in names
        assert "csv_write" in names
        assert "docx_create" in names
        assert "markdown_write" in names
        assert "chart_create" in names

    def test_csv_read_write_roundtrip(self, tmp_path):
        from flowforge.tools.builtin import _make_csv_read_tool, _make_csv_write_tool
        import json

        write_tool = _make_csv_write_tool(tmp_path)
        read_tool = _make_csv_read_tool(tmp_path, 4000)

        data = [
            {"name": "Alice", "age": "30", "city": "Seoul"},
            {"name": "Bob", "age": "25", "city": "Busan"},
        ]
        write_result = write_tool(path="test.csv", data=json.dumps(data))
        assert write_result["ok"] is True
        assert write_result["row_count"] == 2

        read_result = read_tool(path="test.csv")
        assert read_result["ok"] is True
        assert read_result["headers"] == ["name", "age", "city"]
        assert len(read_result["rows"]) == 2
        assert read_result["rows"][0]["name"] == "Alice"

    def test_csv_read_file_not_found(self, tmp_path):
        from flowforge.tools.builtin import _make_csv_read_tool

        tool = _make_csv_read_tool(tmp_path, 4000)
        result = tool(path="nonexistent.csv")
        assert result["ok"] is False
        assert "not found" in result["error"].lower()

    def test_csv_write_rejects_path_outside_project(self, tmp_path):
        from flowforge.tools.builtin import _make_csv_write_tool

        tool = _make_csv_write_tool(tmp_path)
        result = tool(path="/etc/evil.csv", data='[{"a":1}]')
        assert result["ok"] is False
        assert "project root" in result["error"]

    def test_markdown_write(self, tmp_path):
        from flowforge.tools.builtin import _make_markdown_write_tool

        tool = _make_markdown_write_tool(tmp_path)
        content = "# Hello\n\nThis is a test."
        result = tool(path="output/report.md", content=content)
        assert result["ok"] is True
        assert (tmp_path / "output" / "report.md").read_text() == content

    def test_markdown_write_rejects_path_outside_project(self, tmp_path):
        from flowforge.tools.builtin import _make_markdown_write_tool

        tool = _make_markdown_write_tool(tmp_path)
        result = tool(path="/tmp/evil.md", content="bad")
        assert result["ok"] is False
        assert "project root" in result["error"]

    def test_pptx_create_missing_package(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        from unittest.mock import patch
        import builtins

        tool = _make_pptx_create_tool(tmp_path)
        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return original_import(name, *args, **kwargs)

        with patch("builtins.__import__", side_effect=mock_import):
            result = tool(
                path="test.pptx",
                slides='[{"title": "Slide 1", "body": "Hello"}]',
            )
        assert result["ok"] is False
        assert "python-pptx" in result["error"]

    def test_pdf_read_missing_package(self, tmp_path):
        from flowforge.tools.builtin import _make_pdf_read_text_tool
        from unittest.mock import patch
        import builtins

        # Create a dummy file so the file-not-found check passes.
        (tmp_path / "test.pdf").write_bytes(b"%PDF-1.4 dummy")

        tool = _make_pdf_read_text_tool(tmp_path, 4000)
        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pypdf":
                raise ImportError("No module named 'pypdf'")
            return original_import(name, *args, **kwargs)

        with patch("builtins.__import__", side_effect=mock_import):
            result = tool(path="test.pdf")
        assert result["ok"] is False
        assert "pypdf" in result["error"]

    def test_docx_create_missing_package(self, tmp_path):
        from flowforge.tools.builtin import _make_docx_create_tool
        from unittest.mock import patch
        import builtins

        tool = _make_docx_create_tool(tmp_path)
        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "docx":
                raise ImportError("No module named 'docx'")
            return original_import(name, *args, **kwargs)

        with patch("builtins.__import__", side_effect=mock_import):
            result = tool(
                path="test.docx",
                content='[{"type": "paragraph", "text": "Hello"}]',
            )
        assert result["ok"] is False
        assert "python-docx" in result["error"]

    def test_chart_create_missing_package(self, tmp_path):
        from flowforge.tools.builtin import _make_chart_create_tool
        from unittest.mock import patch
        import builtins

        tool = _make_chart_create_tool(tmp_path)
        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "matplotlib":
                raise ImportError("No module named 'matplotlib'")
            return original_import(name, *args, **kwargs)

        with patch("builtins.__import__", side_effect=mock_import):
            result = tool(
                path="chart.png",
                chart_type="bar",
                labels='["A","B","C"]',
                values='[1,2,3]',
            )
        assert result["ok"] is False
        assert "matplotlib" in result["error"]

    def test_chart_create_invalid_type(self, tmp_path):
        from flowforge.tools.builtin import _make_chart_create_tool

        tool = _make_chart_create_tool(tmp_path)
        result = tool(
            path="chart.png",
            chart_type="invalid",
            labels='["A"]',
            values='[1]',
        )
        assert result["ok"] is False
        assert "unsupported" in result["error"].lower()

    def test_pptx_create_invalid_slides_json(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool

        tool = _make_pptx_create_tool(tmp_path)
        result = tool(path="test.pptx", slides="not valid json")
        assert result["ok"] is False
        assert "invalid" in result["error"].lower()

    def test_parse_page_range(self):
        from flowforge.tools.builtin import _parse_page_range

        assert _parse_page_range("1-3", 10) == [0, 1, 2]
        assert _parse_page_range("1,3,5", 10) == [0, 2, 4]
        assert _parse_page_range("2-4,7", 10) == [1, 2, 3, 6]
        # Out of range clamped
        assert _parse_page_range("1-100", 5) == [0, 1, 2, 3, 4]


# ---------------------------------------------------------------------------
# Test: Manifest bridge field
# ---------------------------------------------------------------------------

class TestManifestBridgeField:

    def test_persist_flow_sets_bridge_contract_by_default(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code

        options = DynamicRunOptions(project_root=str(tmp_path))
        record = persist_flow_code(
            flow_name="test_bridge",
            code="# placeholder",
            options=options,
            downstream_flow_route="some_pipeline",
        )
        assert record.bridge == "contract"

    def test_persist_flow_bridge_empty_without_downstream(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code

        options = DynamicRunOptions(project_root=str(tmp_path))
        record = persist_flow_code(
            flow_name="test_no_bridge",
            code="# placeholder",
            options=options,
        )
        assert record.bridge == ""

    def test_persist_flow_bridge_explicit(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code

        options = DynamicRunOptions(project_root=str(tmp_path))
        record = persist_flow_code(
            flow_name="test_shared",
            code="# placeholder",
            options=options,
            bridge="shared_data",
        )
        assert record.bridge == "shared_data"


# ---------------------------------------------------------------------------
# Test: Manifest file locking
# ---------------------------------------------------------------------------

class TestManifestFileLocking:

    def test_lock_file_created_on_persist(self, tmp_path):
        from flowforge.dynamic.manifest import persist_flow_code, resolve_generated_dir

        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(
            flow_name="lock_test",
            code="# lock test",
            options=options,
        )

        lock_path = resolve_generated_dir(options) / "manifest.json.lock"
        assert lock_path.exists()

    def test_concurrent_persist_does_not_corrupt(self, tmp_path):
        """Sequential writes through the lock produce correct manifest."""
        from flowforge.dynamic.manifest import persist_flow_code, load_manifest

        options = DynamicRunOptions(project_root=str(tmp_path))
        persist_flow_code(flow_name="flow_a", code="# a", options=options)
        persist_flow_code(flow_name="flow_b", code="# b", options=options)
        persist_flow_code(flow_name="flow_c", code="# c", options=options)

        manifest = load_manifest(options)
        names = {f.name for f in manifest.flows}
        assert names == {"flow_a", "flow_b", "flow_c"}


# ---------------------------------------------------------------------------
# Test: Meta-flow reduced to 3 steps (no duplicate retry loop)
# ---------------------------------------------------------------------------

class TestMetaFlowConsolidated:

    def test_meta_flow_has_five_steps(self):
        engine = FlowForge.compile(DynamicAgent)
        step_ids = [
            n.id for n in engine.dag.get_all_nodes()
            if n.id.startswith("global._dynamic_generator._generate_and_run.")
            and n.type == NodeType.STEP
        ]
        assert len(step_ids) == 5
        step_id_str = " ".join(step_ids)
        assert "analyse_gap[1]" in step_id_str
        assert "plan_workflow[2]" in step_id_str
        assert "select_capability[3]" in step_id_str
        assert "mcp_provision[4]" in step_id_str
        assert "synthesise_and_inject[5]" in step_id_str


# ---------------------------------------------------------------------------
# Test: Codegen tool catalog includes builtin tools with parameters
# ---------------------------------------------------------------------------

class TestCodegenToolCatalog:

    def test_tool_catalog_includes_builtin_params(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(project_root=str(tmp_path))
        tools = create_builtin_tool_pack(options)

        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            tool_configs=tools,
            dynamic_options=options,
        )

        catalog = gen._format_tool_catalog()

        # Should contain tool names
        assert "web_fetch_url" in catalog
        assert "pptx_create" in catalog
        assert "csv_write" in catalog
        assert "pdf_read_text" in catalog
        assert "chart_create" in catalog

        # Should contain parameter info
        assert "Parameters:" in catalog
        assert "ctx.call_tool" in catalog

    def test_compile_flow_code_applies_dynamic_step_timeout_default(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            generated_step_timeout_seconds=240,
        )
        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            dynamic_options=options,
        )

        meta = gen.compile_flow_code("""
from flowforge import flow, task, step

@flow(name="timeout_flow", prompt="timeout")
class TimeoutFlow:
    @task(name="main", prompt="main")
    class Main:
        @step(order=1, prompt="implicit timeout")
        async def implicit(ctx):
            return {"ok": True}

        @step(order=2, prompt="explicit timeout", timeout_seconds=500)
        async def explicit(ctx):
            return {"ok": True}
""")

        assert meta.tasks[0].steps[0].timeout_seconds == 240
        assert meta.tasks[0].steps[1].timeout_seconds == 500

    def test_tool_catalog_shows_call_example(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(project_root=str(tmp_path))
        tools = create_builtin_tool_pack(options)

        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            tool_configs=tools,
            dynamic_options=options,
        )

        catalog = gen._format_tool_catalog()

        # pptx_create should show path=... and slides=... as required params
        assert 'call_tool("pptx_create"' in catalog
        assert "path=..." in catalog
        assert "slides=..." in catalog

    def test_tool_catalog_shows_agent_skill_call_example(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        skill_dir = tmp_path / "skills" / "clone-coding"
        skill_dir.mkdir(parents=True)
        (skill_dir / "SKILL.md").write_text(
            "---\n"
            "name: clone-coding\n"
            "description: Clone-coding guidance.\n"
            "---\n"
            "# Clone Coding\n",
            encoding="utf-8",
        )

        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            tool_configs=[AgentSkill(path=str(skill_dir))],
        )

        catalog = gen._format_tool_catalog()

        assert "clone-coding (agent-skill)" in catalog
        assert 'ctx.call_llm("instruction <clone-coding>")' in catalog

    def test_tool_catalog_compacts_to_relevant_tools(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            codegen_tool_catalog_max_tools=4,
        )
        tools = create_builtin_tool_pack(options)
        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            tool_configs=tools,
            dynamic_options=options,
        )

        catalog = gen._format_tool_catalog(
            user_query={
                "request": "create a ppt deck",
                "required_tools": ["pptx_create"],
            },
            artifacts=[{"tool": "pptx_create"}],
        )

        assert "pptx_create" in catalog
        assert "lower-relevance tool(s) omitted" in catalog
        assert catalog.count("\n- ") <= 5

    def test_dynamic_policy_lists_declared_mcp_servers(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            mcp_server_commands={"playwright": ["npx", "-y", "@playwright/mcp"]},
            mcp_server_urls={"playwright": "http://localhost:3847/mcp"},
            mcp_server_tools={"playwright": ["browser_navigate"]},
        )
        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            dynamic_options=options,
        )

        policy = gen._format_dynamic_policy()

        assert "playwright" in policy
        assert "browser_navigate" in policy
        assert "codegen_tool_catalog_max_tools" in policy

    def test_playwright_mcp_options_helper_sets_runtime_defaults(self, tmp_path):
        options = DynamicRunOptions.for_playwright_mcp(
            project_root=str(tmp_path),
            port=9999,
        )

        assert options.project_root == str(tmp_path)
        assert options.mcp_server_commands["playwright"] == [
            "npx",
            "-y",
            "@playwright/mcp@latest",
            "--port",
            "9999",
        ]
        assert options.mcp_server_urls["playwright"] == "http://localhost:9999/mcp"
        assert "browser_navigate" in options.mcp_server_tools["playwright"]
        assert "install_dependency" in options.allowed_shell_modes
        assert options.codegen_tool_catalog_max_tools == 8

    def test_figma_mcp_options_helper_sets_remote_defaults(self, tmp_path):
        options = DynamicRunOptions.for_figma_mcp(
            project_root=str(tmp_path),
            authorization="Bearer test-token",
        )

        assert options.project_root == str(tmp_path)
        assert options.mcp_server_urls["figma"] == "https://mcp.figma.com/mcp"
        assert "get_design_context" in options.mcp_server_tools["figma"]
        assert options.mcp_server_headers["figma"] == {
            "Authorization": "Bearer test-token",
        }
        assert options.allowed_shell_modes == ["readonly", "project_exec"]
        assert options.codegen_tool_catalog_max_chars == 4500

    def test_planner_promotes_requirement_gap_metadata(self):
        from flowforge.planner.llm_planner import LLMPlanner

        data = {
            "routes": [],
            "gap_detected": True,
            "reason": "",
            "suggested_flow_name": "",
            "suggested_flow_prompt": "",
            "requirements": [
                {
                    "description": "Need clone coding",
                    "covered": False,
                    "needs_flow": True,
                    "suggested_flow_name": "clone_coding",
                    "suggested_flow_prompt": "Create a clone coding frontend.",
                }
            ],
        }

        promoted = LLMPlanner._promote_gap_from_requirements(data)

        assert promoted["suggested_flow_name"] == "clone_coding"
        assert promoted["suggested_flow_prompt"] == "Create a clone coding frontend."

    def test_planner_infers_gap_from_requirement_metadata(self):
        from flowforge.planner.llm_planner import LLMPlanner

        data = {
            "routes": ["paper_report_pipeline"],
            "gap_detected": False,
            "reason": "",
            "suggested_flow_name": "",
            "suggested_flow_prompt": "",
            "requirements": [
                {
                    "description": "Fetch arXiv papers before reporting",
                    "covered": False,
                    "needs_flow": True,
                    "suggested_flow_name": "arxiv_paper_fetcher",
                    "suggested_flow_prompt": "Fetch recent arXiv AI papers.",
                    "downstream_flow_route": "paper_report_pipeline",
                },
                {
                    "description": "Create the report deck",
                    "covered": True,
                    "matched_route": "paper_report_pipeline",
                },
            ],
        }

        promoted = LLMPlanner._promote_gap_from_requirements(data)

        assert promoted["gap_detected"] is True
        assert promoted["suggested_flow_name"] == "arxiv_paper_fetcher"
        assert promoted["downstream_flow_route"] == "paper_report_pipeline"

    def test_prompt_builder_marks_empty_route_tree(self):
        from flowforge.planner.prompt_builder import PromptBuilder

        engine = FlowForge.compile(EmptyDynamicAgent)
        prompt = PromptBuilder().build(
            "Create a clone coding frontend",
            engine.dag,
            {},
            global_prompt="empty dynamic agent",
        )

        assert "(no user-defined flows available)" in prompt
        assert "fill `suggested_flow_name` plus `suggested_flow_prompt`" in prompt

    def test_strip_markdown_fences_extracts_prefaced_code(self):
        from flowforge.dynamic.generator import _strip_markdown_fences

        text = "Here is the code:\n```python\nfrom flowforge import flow\n```"

        assert _strip_markdown_fences(text) == "from flowforge import flow"

    def test_builtin_shell_tools_accept_string_timeout(self, tmp_path):
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(
            project_root=str(tmp_path),
            allowed_shell_modes=["readonly"],
        )
        tool = next(
            item for item in create_builtin_tool_pack(options)
            if item.name == "shell_readonly"
        )

        result = tool.func("pwd", timeout_seconds="5")

        assert result["ok"] is True
        assert result["exit_code"] == result["returncode"] == 0

    def test_build_context_classifies_document_tools(self, tmp_path):
        from flowforge.dynamic.generator import DynamicFlowGenerator
        from flowforge.tools.builtin import create_builtin_tool_pack

        options = DynamicRunOptions(project_root=str(tmp_path))
        tools = create_builtin_tool_pack(options)

        gen = DynamicFlowGenerator(
            llm_config=LLMConfig(model="test"),
            dag=FlowForge.compile(DynamicAgent).dag,
            tool_configs=tools,
            dynamic_options=options,
        )

        context = gen.build_code_generation_context(
            flow_name="test_flow",
            flow_prompt="test prompt",
            user_query="test query",
        )

        assert "Document tools" in context
        assert "pptx_create" in context
        assert "csv_write" in context
        assert "Utility tools" in context
        assert "web_fetch_url" in context


# ---------------------------------------------------------------------------
# Test: Output artifact detection
# ---------------------------------------------------------------------------

class TestOutputArtifactDetection:

    def test_detect_pptx_from_korean(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("논문 보고서를 PPT로 만들어줘")
        assert len(result) == 1
        assert result[0]["tool"] == "pptx_create"
        assert result[0]["extension"] == ".pptx"

    def test_detect_pptx_from_english(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("create a presentation about AI trends")
        assert len(result) == 1
        assert result[0]["tool"] == "pptx_create"

    def test_detect_csv(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("데이터를 CSV 파일로 저장해줘")
        assert len(result) == 1
        assert result[0]["tool"] == "csv_write"

    def test_detect_chart(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("매출 데이터를 차트로 시각화해줘")
        assert len(result) == 1
        assert result[0]["tool"] == "chart_create"

    def test_detect_multiple_artifacts(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts(
            "논문 요약을 PPT 발표 자료로 만들고 차트도 포함해줘"
        )
        tools = {a["tool"] for a in result}
        assert "pptx_create" in tools
        assert "chart_create" in tools

    def test_detect_docx(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("결과를 워드 문서로 작성해줘")
        assert len(result) == 1
        assert result[0]["tool"] == "docx_create"

    def test_detect_markdown(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("마크다운 보고서를 생성해줘")
        assert len(result) == 1
        assert result[0]["tool"] == "markdown_write"

    def test_detect_nothing_when_no_file_intent(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("최근 AI 논문을 검색해줘")
        assert result == []

    def test_respects_available_tools_filter(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts(
            "PPT로 만들어줘",
            available_tools=["web_fetch_url"],  # pptx_create not available
        )
        assert result == []

    def test_no_duplicate_detection(self):
        from flowforge.dynamic.generator import detect_output_artifacts

        result = detect_output_artifacts("PPT 슬라이드 발표 자료 프레젠테이션")
        assert len(result) == 1  # All map to pptx_create

    def test_format_artifact_instructions(self):
        from flowforge.dynamic.generator import (
            detect_output_artifacts,
            _format_artifact_instructions,
        )

        artifacts = detect_output_artifacts("PPT로 만들어줘")
        instructions = _format_artifact_instructions(artifacts)
        assert "pptx_create" in instructions
        assert "ctx.call_tool" in instructions
        assert "FINAL step" in instructions

    def test_format_empty_artifacts(self):
        from flowforge.dynamic.generator import _format_artifact_instructions

        assert _format_artifact_instructions([]) == ""


# ---------------------------------------------------------------------------
# Test: Enhanced pptx_create tool
# ---------------------------------------------------------------------------

class TestEnhancedPptxCreate:

    def test_cover_layout(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [{"layout": "cover", "title": "Title", "subtitle": "Subtitle"}]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True

    def test_table_layout(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [{
            "layout": "table",
            "title": "Data Table",
            "table": {
                "headers": ["Name", "Age", "City"],
                "rows": [["Alice", "30", "Seoul"], ["Bob", "25", "Busan"]],
            },
        }]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(tmp_path / "test.pptx"))
        slide = prs.slides[0]
        has_table = any(
            hasattr(sh, "has_table") and sh.has_table for sh in slide.shapes
        )
        assert has_table

    def test_comparison_layout(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [{
            "layout": "comparison",
            "title": "A vs B",
            "left": {"heading": "Option A", "bullets": ["Pro 1", "Pro 2"]},
            "right": {"heading": "Option B", "bullets": ["Pro 1", "Pro 2"]},
        }]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True

    def test_dark_theme(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [
            {"layout": "cover", "title": "Dark", "subtitle": "Theme"},
            {"layout": "content", "title": "Content", "bullets": ["A", "B"]},
        ]
        result = tool(path="test.pptx", slides=json.dumps(slides), theme="dark")
        assert result["ok"] is True

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(str(tmp_path / "test.pptx"))
        bg_fill = prs.slides[0].background.fill
        assert bg_fill.fore_color.rgb == RGBColor(0x1E, 0x1E, 0x2E)

    def test_speaker_notes(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [{
            "title": "Notes Test",
            "bullets": ["Point 1"],
            "speaker_note": "Remember to emphasize this.",
        }]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(tmp_path / "test.pptx"))
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "Remember to emphasize" in notes

    def test_notes_alias_and_missing_image_are_tolerated(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [{
            "title": "Optional Image",
            "body": "This should still render.",
            "image_path": "missing.png",
            "notes": "Legacy note alias.",
            "accent": "not-a-hex-color",
        }]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True

        from pptx import Presentation
        prs = Presentation(str(tmp_path / "test.pptx"))
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "Legacy note alias" in notes

    def test_mixed_layouts(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [
            {"layout": "cover", "title": "Report", "subtitle": "Q1 2026"},
            {"layout": "content", "title": "Summary", "bullets": ["A", "B"]},
            {"layout": "table", "title": "Data", "table": {
                "headers": ["X", "Y"], "rows": [["1", "2"]],
            }},
            {"layout": "section", "title": "Details"},
            {"layout": "content", "title": "End", "body": "Thank you"},
        ]
        result = tool(path="test.pptx", slides=json.dumps(slides))
        assert result["ok"] is True
        assert result["slide_count"] == 5
        assert result["native_objects"] is True

    def test_metric_timeline_and_chart_layouts(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        tool = _make_pptx_create_tool(tmp_path)
        slides = [
            {
                "layout": "metric",
                "title": "Signals",
                "metrics": [
                    {"value": "3", "label": "papers"},
                    {"value": "7", "label": "slides"},
                ],
            },
            {
                "layout": "timeline",
                "title": "Workflow",
                "items": [
                    {"label": "1", "title": "Fetch"},
                    {"label": "2", "title": "Summarise"},
                    {"label": "3", "title": "Render"},
                ],
            },
            {
                "layout": "chart",
                "title": "Counts",
                "chart": {
                    "type": "column",
                    "categories": ["A", "B"],
                    "series": [{"name": "Values", "values": [1, 2]}],
                },
            },
        ]
        result = tool(path="rich.pptx", slides=json.dumps(slides), theme="tech")
        assert result["ok"] is True
        assert result["layouts"] == ["metric", "timeline", "chart"]

        from pptx import Presentation
        prs = Presentation(str(tmp_path / "rich.pptx"))
        assert len(prs.slides) == 3
        assert any(shape.has_chart for shape in prs.slides[2].shapes)

    def test_pptx_create_with_symlinked_project_root(self, tmp_path):
        from flowforge.tools.builtin import _make_pptx_create_tool
        import json

        real_root = tmp_path / "real"
        real_root.mkdir()
        linked_root = tmp_path / "linked"
        linked_root.symlink_to(real_root, target_is_directory=True)

        tool = _make_pptx_create_tool(linked_root)
        result = tool(
            path="deck.pptx",
            slides=json.dumps([{"layout": "cover", "title": "Linked root"}]),
        )

        assert result["ok"] is True
        assert (real_root / "deck.pptx").exists()


# ---------------------------------------------------------------------------
# Image creation tool tests
# ---------------------------------------------------------------------------

class TestImageCreateTool:
    """Tests for the image_create builtin tool."""

    def test_pillow_canvas_basic(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        tool = _make_image_create_tool(tmp_path)

        result = tool(
            path="test.png",
            width=200,
            height=100,
            background="#FF0000",
        )
        assert result["ok"] is True
        assert result["width"] == 200
        assert result["height"] == 100
        assert (tmp_path / "test.png").exists()

    def test_pillow_with_elements(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        import json
        tool = _make_image_create_tool(tmp_path)

        elements = json.dumps([
            {"type": "rectangle", "x": 10, "y": 10, "w": 50, "h": 50,
             "color": "#0000FF", "fill": "#00FF00"},
            {"type": "ellipse", "x": 70, "y": 10, "w": 40, "h": 40,
             "color": "#FF0000"},
            {"type": "line", "x1": 0, "y1": 0, "x2": 100, "y2": 100,
             "color": "#000000", "width": 3},
            {"type": "text", "text": "Hello", "x": 10, "y": 60,
             "size": 16, "color": "#333333"},
        ])

        result = tool(path="drawn.png", elements=elements)
        assert result["ok"] is True
        assert result["element_count"] == 4

    def test_jpeg_output(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        tool = _make_image_create_tool(tmp_path)

        result = tool(path="test.jpg", width=100, height=100)
        assert result["ok"] is True
        assert (tmp_path / "test.jpg").exists()

    def test_path_outside_root_rejected(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        tool = _make_image_create_tool(tmp_path)

        result = tool(path="/etc/evil.png")
        assert result["ok"] is False
        assert "project root" in result["error"]

    def test_ai_prompt_without_key(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        import os
        tool = _make_image_create_tool(tmp_path)

        # Ensure no API key is set
        old_key = os.environ.pop("OPENAI_API_KEY", None)
        try:
            result = tool(path="ai.png", ai_prompt="a cat")
            assert result["ok"] is False
            assert "OPENAI_API_KEY" in result["error"]
        finally:
            if old_key:
                os.environ["OPENAI_API_KEY"] = old_key

    def test_invalid_elements_json(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        tool = _make_image_create_tool(tmp_path)

        result = tool(path="bad.png", elements="not json")
        assert result["ok"] is False
        assert "Invalid elements JSON" in result["error"]

    def test_pillow_missing(self, tmp_path):
        from flowforge.tools.builtin import _make_image_create_tool
        from unittest.mock import patch as _patch
        import builtins

        tool = _make_image_create_tool(tmp_path)
        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "PIL" or name.startswith("PIL."):
                raise ImportError("no PIL")
            return original_import(name, *args, **kwargs)

        with _patch("builtins.__import__", side_effect=mock_import):
            result = tool(path="no_pil.png")
            assert result["ok"] is False
            assert "Pillow" in result["error"]


class TestImageArtifactDetection:
    """image_create is detected via artifact rules."""

    def test_detect_image_korean(self):
        from flowforge.dynamic.generator import detect_output_artifacts
        arts = detect_output_artifacts("이미지를 만들어줘")
        tool_names = [a["tool"] for a in arts]
        assert "image_create" in tool_names

    def test_detect_image_english(self):
        from flowforge.dynamic.generator import detect_output_artifacts
        arts = detect_output_artifacts("generate a banner image")
        tool_names = [a["tool"] for a in arts]
        assert "image_create" in tool_names
